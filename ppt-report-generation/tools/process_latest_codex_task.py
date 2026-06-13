from __future__ import annotations

import datetime as dt
import importlib.util
import json
import os
import shutil
import tempfile
import xml.etree.ElementTree as ET
from zipfile import ZIP_DEFLATED, ZipFile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


ROOT = Path.cwd()
TOOL_PATH = next(ROOT.glob("PDF*/pdf_ppt_tool.py"))
LATEST_TASK_POINTER = next(ROOT.glob("PDF*/*Codex*.txt"))
LATEST_TASK_PATH = Path(LATEST_TASK_POINTER.read_text(encoding="utf-8").strip())

spec = importlib.util.spec_from_file_location("ppttool", TOOL_PATH)
ppttool = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(ppttool)

TASK = json.loads(LATEST_TASK_PATH.read_text(encoding="utf-8"))
TEMPLATE = Path(TASK["template_path"])
SOURCE = Path(TASK["source_dir"])
OUTPUT_DIR = Path(TASK["output_dir"])
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

STAMP = dt.datetime.now().strftime("%Y%m%d-%H%M%S")
BASE_NAME = f"{TASK['title']} - Codex处理版-{STAMP}"
PPTX_OUT = OUTPUT_DIR / f"{BASE_NAME}.pptx"
PLAN_OUT = OUTPUT_DIR / f"{BASE_NAME}_codex_plan.json"

TITLE_COLOR = RGBColor(31, 41, 55)
BODY_COLOR = RGBColor(31, 41, 55)
PANEL = RGBColor(246, 248, 251)
LINE = RGBColor(195, 202, 214)


def nonempty_text_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
            shapes.append(shape)
    return sorted(shapes, key=lambda item: (int(item.top), int(item.left)))


def first_style(shape, fallback_size=24):
    style = {
        "name": "Microsoft YaHei",
        "size": fallback_size,
        "bold": False,
        "italic": False,
        "color": BODY_COLOR,
        "align": None,
    }
    try:
        style["align"] = shape.text_frame.paragraphs[0].alignment
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                names = ppttool.run_font_names(run)
                if names:
                    style["name"] = next((name for kind, name in names if kind == "east_asian"), names[0][1])
                if run.font.size:
                    style["size"] = int(run.font.size.pt)
                if run.font.bold is not None:
                    style["bold"] = bool(run.font.bold)
                if run.font.italic is not None:
                    style["italic"] = bool(run.font.italic)
                try:
                    if run.font.color.rgb:
                        style["color"] = run.font.color.rgb
                except Exception:
                    pass
                if run.text.strip():
                    return style
    except Exception:
        pass
    return style


def set_text(shape, text, size=None, bold=None, color=None):
    style = first_style(shape, fallback_size=size or 24)
    if size is not None:
        style["size"] = size
    if bold is not None:
        style["bold"] = bool(bold)
    if color is not None:
        style["color"] = color

    lines = str(text).split("\n") if text is not None else [""]
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        paragraph.text = line
        if style["align"] is not None:
            paragraph.alignment = style["align"]
        for run in paragraph.runs:
            ppttool.set_run_font(run, style["name"])
            run.font.size = Pt(style["size"])
            run.font.bold = style["bold"]
            run.font.italic = style["italic"]
            try:
                run.font.color.rgb = style["color"]
            except Exception:
                pass


def replace_slide_text(prs, slide_index, texts, sizes=None, bolds=None):
    slide = prs.slides[slide_index - 1]
    shapes = nonempty_text_shapes(slide)
    for index, text in enumerate(texts):
        if index >= len(shapes):
            break
        set_text(shapes[index], text)
    return slide


def clear_pictures(slide):
    slide_w = max(int(shape.left) + int(shape.width) for shape in slide.shapes)
    slide_h = max(int(shape.top) + int(shape.height) for shape in slide.shapes)
    for shape in list(slide.shapes):
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            if ppttool.is_full_slide_shape(shape, slide_w, slide_h):
                continue
            shape.element.getparent().remove(shape.element)


def add_panel(slide, left, top, width, height, title, items, title_size=19, body_size=15):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL
    rect.line.color.rgb = LINE

    box = slide.shapes.add_textbox(Emu(left + 220000), Emu(top + 170000), Emu(width - 440000), Emu(height - 300000))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    paragraph = frame.paragraphs[0]
    paragraph.text = title
    for run in paragraph.runs:
        ppttool.set_run_font(run, "微软雅黑")
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR

    for item in items:
        paragraph = frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(4)
        for run in paragraph.runs:
            ppttool.set_run_font(run, "微软雅黑")
            run.font.size = Pt(body_size)
            run.font.color.rgb = BODY_COLOR


def media_item(path):
    path = Path(path)
    return {"path": str(path), "name": path.name, "relative": path.name}


def media_label(path):
    kind = ppttool.media_kind_from_path(Path(path))
    return "音频展示类：" if kind == "audio" else "产出展示类："


def set_shape_link(shape, path):
    try:
        shape.click_action.hyperlink.address = str(path)
    except Exception:
        pass


def add_media_showcase(slide, slide_number, left, top, width, height, title, paths):
    media = [media_item(path) for path in paths if Path(path).exists()]
    if not media:
        return []
    item = media[0]
    shapes = nonempty_text_shapes(slide)
    target = shapes[-1] if shapes else slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    set_text(target, f"{media_label(item['path'])}\n1.{item['name']}")
    set_shape_link(target, item["path"])
    return [
        {
            "slide": slide_number,
            "path": item["path"],
            "name": item["name"],
            "kind": ppttool.media_kind_from_path(Path(item["path"])),
            "mode": "single_link_text",
        }
    ]


def move_last_slide_before(prs, before_index):
    sld_id = prs.slides._sldIdLst[-1]
    del prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.insert(before_index, sld_id)


def append_single_media_pages(prs, slides, source_slide_index, project_title, paths):
    insertions = []
    for path in paths:
        path = Path(path)
        if not path.exists():
            continue
        new_slide = ppttool.duplicate_slide(prs, source_slide_index - 1)
        before_index = len(prs.slides) - 2
        move_last_slide_before(prs, before_index)
        slide_number = before_index + 1
        shapes = nonempty_text_shapes(new_slide)
        if len(shapes) >= 2:
            set_text(shapes[0], project_title)
        insertions.extend(add_media_showcase(new_slide, slide_number, 700000, 1500000, 10750000, 3600000, project_title, [path]))
        slides.append({"slide": slide_number, "title": f"{project_title} / {path.name}"})
    return insertions


def add_picture_fit(slide, image_path, left, top, width, height):
    path = Path(image_path)
    if not path.exists():
        return False
    with Image.open(path) as image:
        image_width, image_height = image.size
    scale = min(width / image_width, height / image_height)
    display_width = int(image_width * scale)
    display_height = int(image_height * scale)
    display_left = int(left + (width - display_width) / 2)
    display_top = int(top + (height - display_height) / 2)
    slide.shapes.add_picture(
        str(path),
        Emu(display_left),
        Emu(display_top),
        width=Emu(display_width),
        height=Emu(display_height),
    )
    return True


MEDIA_EXTS = {".mp4", ".mov", ".m4v", ".avi", ".wmv", ".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}


def is_av_media_name(name):
    return Path(str(name).split("#", 1)[0].split("?", 1)[0]).suffix.lower() in MEDIA_EXTS


def remove_xml_elements_by_localname(root, local_names):
    changed = False
    parent_map = {child: parent for parent in root.iter() for child in list(parent)}
    for element in list(root.iter()):
        local_name = str(element.tag).rsplit("}", 1)[-1]
        if local_name in local_names:
            parent = parent_map.get(element)
            if parent is not None:
                parent.remove(element)
                changed = True
    return changed


def strip_embedded_av_media(pptx_path):
    removed = []
    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    tmp_path = Path(tmp_name)
    try:
        with ZipFile(pptx_path, "r") as source_zip, ZipFile(tmp_path, "w", ZIP_DEFLATED) as output_zip:
            for item in source_zip.infolist():
                name = item.filename
                data = source_zip.read(name)

                if name.startswith("ppt/media/") and is_av_media_name(name):
                    removed.append(name)
                    continue

                if name.endswith(".rels"):
                    root = ET.fromstring(data)
                    changed = False
                    for rel in list(root):
                        target = rel.attrib.get("Target", "")
                        rel_type = rel.attrib.get("Type", "")
                        if ("media/" in target and is_av_media_name(target)) or rel_type.endswith(("/audio", "/video", "/media")):
                            root.remove(rel)
                            changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                elif name.endswith(".xml") and b"sndAc" in data:
                    root = ET.fromstring(data)
                    if remove_xml_elements_by_localname(root, {"sndAc"}):
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                elif name == "[Content_Types].xml":
                    root = ET.fromstring(data)
                    changed = False
                    for child in list(root):
                        part_name = child.attrib.get("PartName", "")
                        extension = child.attrib.get("Extension", "")
                        if (part_name.startswith("/ppt/media/") and is_av_media_name(part_name)) or (
                            extension and f".{extension.lower()}" in MEDIA_EXTS
                        ):
                            root.remove(child)
                            changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)

                output_zip.writestr(item, data)
        shutil.move(str(tmp_path), str(pptx_path))
    finally:
        if tmp_path.exists():
            try:
                tmp_path.unlink()
            except Exception:
                pass
    return removed


def shape_type_value(shape):
    try:
        return int(getattr(shape, "shape_type", 0))
    except Exception:
        return None


def is_media_shape(shape):
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.MEDIA or shape_type_value(shape) == 16


def rel_ids_from_shape(shape):
    rel_ids = set()
    for element in shape.element.iter():
        for value in element.attrib.values():
            if isinstance(value, str) and value.startswith("rId"):
                rel_ids.add(value)
    return rel_ids


def remove_shape_with_relationships(slide, shape):
    for rel_id in rel_ids_from_shape(shape):
        try:
            slide.part.drop_rel(rel_id)
        except Exception:
            try:
                del slide.part.rels._rels[rel_id]
            except Exception:
                pass
    parent = shape.element.getparent()
    if parent is not None:
        parent.remove(shape.element)


def capture_media_rects(prs):
    rects = {}
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if is_media_shape(shape):
                rects.setdefault(
                    slide_index,
                    (int(shape.left), int(shape.top), int(shape.width), int(shape.height)),
                )
                break
    return rects


def clean_old_media_shapes(prs):
    removed = 0
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if is_media_shape(shape):
                remove_shape_with_relationships(slide, shape)
                removed += 1
    return removed


def slide_index(prs, target_slide):
    for index, slide in enumerate(prs.slides):
        if slide == target_slide:
            return index
    raise ValueError("slide is not in presentation")


def slide_number(prs, target_slide):
    return slide_index(prs, target_slide) + 1


def remove_slide(prs, target_slide):
    index = slide_index(prs, target_slide)
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    del prs.slides._sldIdLst[index]
    try:
        prs.part.drop_rel(rel_id)
    except Exception:
        pass


def move_last_slide_after(prs, after_index):
    slide_id = prs.slides._sldIdLst[-1]
    del prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.insert(after_index + 1, slide_id)


def duplicate_slide_after(prs, source_slide, after_slide):
    source_index = slide_index(prs, source_slide)
    after_index = slide_index(prs, after_slide)
    new_slide = ppttool.duplicate_slide(prs, source_index)
    move_last_slide_after(prs, after_index)
    return new_slide


def existing_media(paths):
    return [Path(path) for path in paths if Path(path).exists() and Path(path).suffix.lower() in MEDIA_EXTS]


def media_under(*parts):
    base = SOURCE.joinpath(*parts)
    if not base.exists():
        return []
    paths = [
        path
        for path in base.rglob("*")
        if path.is_file()
        and path.suffix.lower() in MEDIA_EXTS
        and "_internal" not in path.parts
        and ".dist-info" not in path.as_posix()
    ]
    return sorted(paths, key=lambda path: str(path.relative_to(SOURCE)).lower())


def rect_to_inches(rect):
    return [round(value / 914400, 4) for value in rect]


def media_display_label(path):
    return "音频展示" if ppttool.media_kind_from_path(Path(path)) == "audio" else "产出展示"


def default_media_rect(prs):
    return (
        int(prs.slide_width * 0.08),
        int(prs.slide_height * 0.28),
        int(prs.slide_width * 0.84),
        int(prs.slide_height * 0.54),
    )


def media_rect_for_slide(prs, media_rects, slide_number_value):
    return media_rects.get(slide_number_value) or default_media_rect(prs)


def add_clickable_media_placeholder(slide, path, rect):
    path = Path(path)
    kind = ppttool.media_kind_from_path(path)
    left, top, width, height = rect
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    if kind == "audio":
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(248, 250, 252)
        box.line.color.rgb = RGBColor(148, 163, 184)
        color = RGBColor(31, 41, 55)
        text = f"音频\n{path.name}"
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(31, 41, 55)
        box.line.color.rgb = RGBColor(37, 99, 235)
        color = RGBColor(255, 255, 255)
        text = f"视频\n{path.name}"

    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(box, text, size=15, bold=False, color=color)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    set_shape_link(box, path)
    return box


def set_media_slide(prs, slide, project_title, path, rect):
    path = Path(path)
    current_number = slide_number(prs, slide)
    title_shape = nonempty_text_shapes(slide)[0] if nonempty_text_shapes(slide) else None
    if title_shape is None:
        title_shape = slide.shapes.add_textbox(Emu(700000), Emu(250000), Emu(10800000), Emu(600000))
    title_text = f"{project_title} {media_display_label(path)}：{path.name}"
    set_text(title_shape, title_text)
    set_shape_link(title_shape, path)

    rect = rect or default_media_rect(prs)
    add_clickable_media_placeholder(slide, path, rect)
    return {
        "slide": current_number,
        "path": str(path),
        "name": path.name,
        "kind": ppttool.media_kind_from_path(path),
        "mode": "single_media_per_template_slot",
        "rect": rect_to_inches(rect),
    }


def ensure_media_slots(prs, base_slots, paths):
    paths = [Path(path) for path in paths if Path(path).exists()]
    if not paths:
        for slide, _rect in reversed(base_slots):
            remove_slide(prs, slide)
        return []

    slots = list(base_slots)
    source_slide, source_rect = slots[-1]
    after_slide = source_slide
    while len(slots) < len(paths):
        new_slide = duplicate_slide_after(prs, source_slide, after_slide)
        slots.append((new_slide, source_rect))
        after_slide = new_slide

    for unused_slide, _rect in reversed(slots[len(paths) :]):
        remove_slide(prs, unused_slide)
    return slots[: len(paths)]


def build_media_groups():
    return {
        "L": existing_media(
            [
                SOURCE / "L" / "loading.mp4",
                SOURCE / "L" / "UI.mp4",
                SOURCE / "L" / "郑青" / "试听.mp4",
            ]
        ),
        "LL": media_under("LL"),
        "U": media_under("U"),
        "FM": media_under("FM"),
        "EF": media_under("EF"),
        "TA": media_under("音频技术"),
    }


def build_deck():
    prs = Presentation(str(TEMPLATE))
    media_rects = capture_media_rects(prs)
    removed_old_media = clean_old_media_shapes(prs)
    media_groups = build_media_groups()

    media_base_slots = {
        "L": [(prs.slides[5], media_rects.get(6)), (prs.slides[6], media_rects.get(7))],
        "LL": [
            (prs.slides[10], media_rects.get(11)),
            (prs.slides[11], media_rects.get(12)),
            (prs.slides[12], media_rects.get(13)),
            (prs.slides[13], media_rects.get(14)),
        ],
        "U": [(prs.slides[16], media_rects.get(17))],
        "FM": [(prs.slides[20], media_rects.get(21))],
        "EF": [(prs.slides[24], media_rects.get(25)), (prs.slides[25], media_rects.get(26))],
        "TA": [(prs.slides[27], media_rects.get(28)), (prs.slides[28], media_rects.get(29))],
    }

    replace_slide_text(prs, 1, ["音频中台 工作汇报", "（6月1日-6月12日 周报）"])
    replace_slide_text(
        prs,
        2,
        [
            "人员名单",
            "音频中台 | 虞鹏 | 郑青 | 蔡新星 | 安家辉 | 高世佳 | 冯博宇",
            "朱静雅 | 关晗怡 | 吴帅斐 | 吉禹鹏",
        ],
    )

    replace_slide_text(
        prs,
        3,
        [
            "重点工作1——L项目",
            "参与人员\n音效设计：郑青\n音频技术：冯博宇，蔡新星\n本期重点：0618、0716版本对接，音效制作与语音标准化",
        ],
    )
    replace_slide_text(
        prs,
        4,
        ["设计&沟通类：\n1. 0618、0716版本音频需求对接\n2. 语音标准化流程整理\n3. UI / loading 相关音频需求沟通"],
    )
    replace_slide_text(
        prs,
        5,
        ["资源制作类：\n1. 音效制作 28 个\n2. loading、UI 视频素材 2 条\n3. 配合版本联调与资源整理"],
    )

    replace_slide_text(
        prs,
        8,
        [
            "重点工作2——LL项目",
            "参与人员\n音效设计：关晗怡，安家辉，吴帅斐，虞鹏\n音乐设计：朱静雅，吉禹鹏，虞鹏\n音频策划：关晗怡，虞鹏\n音频技术：冯博宇，蔡新星",
        ],
    )
    replace_slide_text(
        prs,
        9,
        [
            "重点工作2——LL项目",
            "设计&沟通类：\n1. AI工具提炼策划、美术与 featurelist 信息\n2. AI日语识别检测优化、剧情自动录屏需求\n3. External Source 编辑器与联调工具需求梳理\n4. 项目周会、功能会与版本跑测跟进",
        ],
    )
    replace_slide_text(
        prs,
        10,
        [
            "重点工作2——LL项目",
            "资源制作 / 问题处理类：\n1. 战斗技能、演出技能与英雄入场音效补齐\n2. 音效缺漏补齐日志与问题定位\n3. Prefab / Spine / Timeline 音频事件排查\n4. AI 音频创作工具演示与对比资料整理",
        ],
    )

    replace_slide_text(
        prs,
        15,
        [
            "重点工作3——U项目",
            "参与人员：\n音效师：郑青\n音乐设计：郑青 / 朱静雅 / 虞鹏\n音频策划：郑青\n音频技术：冯博宇，蔡新星",
        ],
    )
    replace_slide_text(
        prs,
        16,
        ["资源制作类：\n视频 2 条\n音效 24 个\n音乐剪辑混音\n语音发包 / 回包\nWwise 维护与 Unity 挂接"],
    )

    replace_slide_text(
        prs,
        18,
        ["重点工作4  ——FM项目", "参与人员\n音效设计：冯博宇\n音乐设计：冯博宇\n音频策划：冯博宇\n音频技术：冯博宇，蔡新星"],
    )
    replace_slide_text(prs, 19, ["设计&沟通类：\nFM：OB24 音效制作\n多语言分支维护\nFM 项目音频看板更新"])
    replace_slide_text(prs, 20, ["资源制作类：\n1. 技能音效 18 个\n2. 资源入库与分支管理\n3. 项目音频看板维护"])

    replace_slide_text(
        prs,
        22,
        ["重点工作5  ——EF项目", "参与人员\n音效设计：虞鹏\n音乐设计：虞鹏\n音频策划：虞鹏\n音频技术：蔡新星"],
    )
    replace_slide_text(prs, 23, ["EF", "设计&沟通类：\n1. UI 替换音频跟进\n2. Stamina 功能音频验证\n3. 成就系统音频联调"])
    replace_slide_text(prs, 24, ["EF", "资源制作类：\n1. 0610 UI替换.mp4\n2. Stamina.mp4\n3. 成就.mp4"])

    replace_slide_text(
        prs,
        27,
        [
            "TA",
            "平台：开发 WwiseExternalSourceManager；开发工具管理器服务器端管理工具。\nLL：继续开发 UIPrefabFindHelper，更新入库工具，真机测试。\nU：OB69 音频功能接入、TW 版本功能单。\nL：0618、0716 版本对接，语音标准化。\n其他：Unity 自动录制特效、自动特效挂接音频、关卡文案导入、整合看板。",
        ],
    )

    dashboard_image = SOURCE / "音频技术" / "冯博宇大看板ver2分支管理.jpg"
    clear_pictures(prs.slides[29])
    replace_slide_text(
        prs,
        30,
        [
            "其他内容",
            "工具与流程沉淀\n1. 联调工具：实时显示特效、Spine 动画和 Timeline，点击可定位资源文件。\n2. Spine 查错工具：遍历探索 Spine 音频事件及父级衰减配置。\n3. External Source 编辑器：按标签提取资源并编辑音量。\n4. AI 工具制作与游戏音频 AI 创作工具对比。",
        ],
    )
    shapes = nonempty_text_shapes(prs.slides[29])
    if len(shapes) > 1:
        shapes[1].width = Emu(5850000)
    add_picture_fit(prs.slides[29], dashboard_image, 6450000, 1450000, 4700000, 3500000)
    replace_slide_text(prs, 31, ["THANK YOU"])

    media_insertions = []
    section_titles = {
        "L": "重点工作1——L项目",
        "LL": "重点工作2——LL项目",
        "U": "重点工作3——U项目",
        "FM": "重点工作4——FM项目",
        "EF": "重点工作5——EF项目",
        "TA": "TA",
    }
    prepared_slots = {}
    for key in ["L", "LL", "U", "FM", "EF", "TA"]:
        prepared_slots[key] = ensure_media_slots(prs, media_base_slots[key], media_groups[key])

    for key in ["L", "LL", "U", "FM", "EF", "TA"]:
        for (slide, rect), path in zip(prepared_slots[key], media_groups[key]):
            media_insertions.append(set_media_slide(prs, slide, section_titles[key], path, rect))

    slides = []
    for index, slide in enumerate(prs.slides, 1):
        text_shapes = nonempty_text_shapes(slide)
        title = text_shapes[0].text.replace("\n", " ")[:80] if text_shapes else f"Slide {index}"
        slides.append({"slide": index, "title": title})

    slides.append({"meta": "removed_old_template_media_objects", "count": removed_old_media})
    return prs, slides, media_insertions


def main():
    prs, slides, media_insertions = build_deck()
    PLAN_OUT.write_text(
        json.dumps(
            {
                "task": str(LATEST_TASK_PATH),
                "template": str(TEMPLATE),
                "source": str(SOURCE),
                "slides": slides,
                "media_insertions": media_insertions,
                "excluded_noise": ["WwiseExternalSourceManager/_internal", "~$临时文件", ".DS_Store"],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    prs.save(str(PPTX_OUT))
    stripped_media = strip_embedded_av_media(PPTX_OUT)
    if stripped_media:
        Presentation(str(PPTX_OUT)).save(str(PPTX_OUT))
        plan_data = json.loads(PLAN_OUT.read_text(encoding="utf-8"))
        plan_data["stripped_embedded_av_media"] = stripped_media
        plan_data["normalized_after_media_strip"] = True
        PLAN_OUT.write_text(json.dumps(plan_data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"STRIPPED_EMBEDDED_AV={len(stripped_media)}")
    print(f"MEDIA_LINKS={len(media_insertions)}")
    print(f"PPTX={PPTX_OUT}")
    print(f"PLAN={PLAN_OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
