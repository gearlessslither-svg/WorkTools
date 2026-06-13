from __future__ import annotations

import argparse
import copy
import datetime as dt
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import traceback
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Callable
from xml.etree import ElementTree as ET

import fitz
import openpyxl
import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.opc.package import _Relationship
from pptx.util import Emu, Pt

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox
    from tkinter import ttk
except Exception:  # pragma: no cover - CLI mode can still work.
    tk = None
    filedialog = None
    messagebox = None
    ttk = None


APP_NAME = "PPTX 智能生成工具"
EMU_PER_INCH = 914400
POINTS_PER_INCH = 72
DEFAULT_FONT = "Microsoft YaHei"
WORK_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = WORK_DIR.parent
DEFAULT_OUTPUT_DIR = WORK_DIR / "生成结果"
PAIR_PACKAGE_SCRIPT = WORK_DIR / "process_pair_package_ppt.py"

TEXT_EXTS = {".txt", ".md", ".markdown"}
DOCX_EXTS = {".docx"}
XLSX_EXTS = {".xlsx", ".xlsm"}
PDF_EXTS = {".pdf"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"}
MEDIA_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".wav", ".mp3", ".flac", ".ogg", ".m4a"}
VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv"}
AUDIO_EXTS = {".wav", ".mp3", ".flac", ".ogg", ".m4a"}
MAX_EMBED_MEDIA_BYTES = 40 * 1024 * 1024
MAX_POWERPOINT_INSERT_MEDIA_BYTES = 30 * 1024 * 1024
SUPPORTED_TEMPLATE_EXTS = {".pptx"}
MAX_TEXT_CHARS_PER_FILE = 9000
MAX_SOURCE_FILES = 400
MAX_SECTIONS = 40


LogFn = Callable[[str], None]


def noop_log(_: str) -> None:
    return None


def safe_console_print(message: str) -> None:
    try:
        print(message, flush=True)
    except UnicodeEncodeError:
        data = (str(message) + "\n").encode("utf-8", errors="backslashreplace")
        sys.stdout.buffer.write(data)
        sys.stdout.buffer.flush()


def stamp() -> str:
    return dt.datetime.now().strftime("%Y%m%d-%H%M%S")


def safe_filename(text: str, default: str = "生成报告") -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\r\n\t]+', "_", text).strip(" ._")
    return cleaned[:80] or default


def clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def emu_to_inches(value: int | float) -> float:
    return float(value) / EMU_PER_INCH


def inches_to_emu(value: float) -> int:
    return int(value * EMU_PER_INCH)


def emu_to_points(value: int | float) -> float:
    return float(value) / EMU_PER_INCH * POINTS_PER_INCH


def points_to_emu(value: int | float) -> int:
    return int(float(value) / POINTS_PER_INCH * EMU_PER_INCH)


def tuple_to_rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(int(color[0]), int(color[1]), int(color[2]))


def hex_to_tuple(value: str | None, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    if not value:
        return fallback
    value = value.strip().lstrip("#")
    if len(value) != 6:
        return fallback
    try:
        return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
    except ValueError:
        return fallback


def tuple_to_hex(color: tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*color)


def rgb_obj_to_tuple(value: Any) -> tuple[int, int, int] | None:
    if value is None:
        return None
    try:
        return int(value[0]), int(value[1]), int(value[2])
    except Exception:
        return None


def color_brightness(color: tuple[int, int, int]) -> float:
    r, g, b = color
    return (r * 0.299 + g * 0.587 + b * 0.114) / 255


def color_saturation(color: tuple[int, int, int]) -> float:
    r, g, b = [x / 255 for x in color]
    return max(r, g, b) - min(r, g, b)


def readable_text_color(bg: tuple[int, int, int]) -> tuple[int, int, int]:
    return (24, 31, 42) if color_brightness(bg) > 0.58 else (248, 250, 252)


def lighten(color: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    return tuple(int(c + (255 - c) * amount) for c in color)


def darken(color: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    return tuple(int(c * (1 - amount)) for c in color)


def ensure_dir(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path


def rel_path(path: Path, base: Path) -> str:
    try:
        return str(path.relative_to(base))
    except ValueError:
        return str(path)


def common_windows_cjk_font() -> str | None:
    candidates = [
        Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\msyhbd.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ]
    for path in candidates:
        if path.exists():
            return str(path)
    return None


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16", "big5"):
        try:
            return path.read_text(encoding=encoding, errors="strict")
        except Exception:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_docx_text(path: Path) -> str:
    paragraphs: list[str] = []
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as docx:
        try:
            xml = docx.read("word/document.xml")
        except KeyError:
            return ""
    root = ET.fromstring(xml)
    for para in root.findall(".//w:p", ns):
        texts = [node.text or "" for node in para.findall(".//w:t", ns)]
        line = "".join(texts).strip()
        if line:
            paragraphs.append(line)
    return "\n".join(paragraphs)


def extract_pdf_text(path: Path, max_pages: int = 20) -> str:
    chunks: list[str] = []
    with fitz.open(path) as pdf:
        for page_index, page in enumerate(pdf):
            if page_index >= max_pages:
                chunks.append(f"...已省略后续 {len(pdf) - max_pages} 页")
                break
            text = page.get_text().strip()
            if text:
                chunks.append(text)
    return "\n\n".join(chunks)


def extract_xlsx_text(path: Path, max_sheets: int = 4, max_rows: int = 36, max_cols: int = 10) -> str:
    chunks: list[str] = []
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    try:
        for sheet_index, ws in enumerate(wb.worksheets[:max_sheets]):
            rows: list[str] = []
            for row_index, row in enumerate(ws.iter_rows(values_only=True)):
                if row_index >= max_rows:
                    break
                values = []
                for cell in row[:max_cols]:
                    if cell is None:
                        continue
                    value = str(cell).strip()
                    if value:
                        values.append(value)
                if values:
                    rows.append(" | ".join(values))
            if rows:
                chunks.append(f"[{ws.title}]\n" + "\n".join(rows))
    finally:
        wb.close()
    return "\n\n".join(chunks)


def extract_text_from_file(path: Path) -> tuple[str, str | None]:
    suffix = path.suffix.lower()
    try:
        if suffix in TEXT_EXTS:
            return read_text_file(path), None
        if suffix in DOCX_EXTS:
            return extract_docx_text(path), None
        if suffix in PDF_EXTS:
            return extract_pdf_text(path), None
        if suffix in XLSX_EXTS:
            return extract_xlsx_text(path), None
        if suffix == ".doc":
            return "", "旧版 .doc 暂未直接解析，建议另存为 .docx。"
    except Exception as exc:
        return "", f"解析失败：{exc}"
    return "", None


def clean_line(line: str) -> str:
    line = re.sub(r"\s+", " ", line.replace("\u3000", " ")).strip()
    line = re.sub(r"^[•\-–—*·\d.、\)\(]+", "", line).strip()
    return line


def split_meaningful_lines(text: str, max_lines: int = 160) -> list[str]:
    lines: list[str] = []
    for raw in re.split(r"[\r\n]+", text):
        line = clean_line(raw)
        if not line or len(line) < 3:
            continue
        if re.fullmatch(r"[\W_]+", line):
            continue
        lines.append(line[:180])
        if len(lines) >= max_lines:
            break
    return lines


def score_line(line: str) -> int:
    keywords = [
        "完成",
        "优化",
        "新增",
        "修复",
        "支持",
        "接入",
        "上线",
        "验证",
        "排查",
        "问题",
        "风险",
        "计划",
        "下周",
        "进度",
        "总结",
        "设计",
        "工具",
        "需求",
        "联调",
        "测试",
        "交付",
        "产出",
    ]
    score = 0
    for kw in keywords:
        if kw in line:
            score += 3
    if 8 <= len(line) <= 80:
        score += 2
    if any(ch in line for ch in "：:"):
        score += 1
    return score


def summarize_lines(text: str, limit: int = 7) -> list[str]:
    lines = split_meaningful_lines(text)
    if not lines:
        return []
    deduped: list[str] = []
    seen: set[str] = set()
    for line in lines:
        key = re.sub(r"\W+", "", line.lower())[:60]
        if key and key not in seen:
            seen.add(key)
            deduped.append(line)
    ranked = sorted(enumerate(deduped), key=lambda item: (-score_line(item[1]), item[0]))
    selected = sorted(ranked[:limit], key=lambda item: item[0])
    return [line for _, line in selected]


def file_kind(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTS:
        return "image"
    if suffix in MEDIA_EXTS:
        return "media"
    if suffix in TEXT_EXTS | DOCX_EXTS | XLSX_EXTS | PDF_EXTS or suffix == ".doc":
        return "document"
    return "other"


def should_skip(path: Path) -> bool:
    parts = {part.lower() for part in path.parts}
    name = path.name.lower()
    ignored = {
        ".git",
        ".codex",
        ".venv",
        "venv",
        "_internal",
        "__macosx",
        "__pycache__",
        "node_modules",
        "生成结果",
        "pdf生成工具",
    }
    ignored_names = {".ds_store", "thumbs.db", "desktop.ini"}
    if name in ignored_names or path.name.startswith("~$"):
        return True
    if any(part.endswith(".dist-info") or part.endswith(".egg-info") for part in parts):
        return True
    return bool(parts & ignored)


def scan_source_folder(source_dir: Path, log: LogFn = noop_log) -> dict[str, Any]:
    source_dir = source_dir.resolve()
    if not source_dir.exists() or not source_dir.is_dir():
        raise FileNotFoundError(f"资料文件夹不存在：{source_dir}")

    log(f"读取资料文件夹：{source_dir}")
    files = sorted(
        [p for p in source_dir.rglob("*") if p.is_file() and not should_skip(p)],
        key=lambda item: str(item).lower(),
    )
    files = files[:MAX_SOURCE_FILES]

    grouped: dict[str, list[Path]] = defaultdict(list)
    for path in files:
        relative = path.relative_to(source_dir)
        group = relative.parts[0] if len(relative.parts) > 1 else "根目录"
        grouped[group].append(path)

    sections: list[dict[str, Any]] = []
    warnings: list[str] = []
    total_text_chars = 0
    total_images = 0
    total_media = 0
    total_docs = 0

    for group_name, group_files in grouped.items():
        docs: list[dict[str, Any]] = []
        images: list[dict[str, Any]] = []
        media: list[dict[str, Any]] = []
        others: list[str] = []

        for path in sorted(group_files, key=lambda x: str(x).lower()):
            kind = file_kind(path)
            if kind == "image":
                total_images += 1
                images.append({"path": str(path), "name": path.name, "relative": rel_path(path, source_dir)})
            elif kind == "media":
                total_media += 1
                media.append({"path": str(path), "name": path.name, "relative": rel_path(path, source_dir)})
            elif kind == "document":
                total_docs += 1
                text, warning = extract_text_from_file(path)
                if warning:
                    warnings.append(f"{rel_path(path, source_dir)}：{warning}")
                text = text[:MAX_TEXT_CHARS_PER_FILE]
                total_text_chars += len(text)
                docs.append(
                    {
                        "path": str(path),
                        "name": path.name,
                        "relative": rel_path(path, source_dir),
                        "text": text,
                        "summary": summarize_lines(text, 5),
                    }
                )
            else:
                others.append(rel_path(path, source_dir))

        if not docs and not images and not media:
            continue

        combined_text = "\n".join(doc["text"] for doc in docs if doc["text"])
        bullets = summarize_lines(combined_text, 7)
        if not bullets:
            if images:
                bullets.append(f"包含 {len(images)} 张图片素材，可用于页面配图。")
            if media:
                bullets.append(f"包含 {len(media)} 个音视频文件，已整理到媒体清单。")
            if docs and not combined_text.strip():
                bullets.append("包含文档文件，但未提取到足够正文。")
        if not bullets:
            bullets.append("该模块暂无可解析文本，已保留文件清单供人工补充。")

        sections.append(
            {
                "title": group_name,
                "bullets": bullets,
                "documents": docs,
                "images": images,
                "media": media,
                "other_files": others,
            }
        )

    sections = sorted(sections, key=lambda s: (s["title"] == "根目录", s["title"]))
    if len(sections) > MAX_SECTIONS:
        warnings.append(f"资料分组超过 {MAX_SECTIONS} 个，仅生成前 {MAX_SECTIONS} 个分组。")
        sections = sections[:MAX_SECTIONS]

    inventory = {
        "source_dir": str(source_dir),
        "sections": sections,
        "warnings": warnings,
        "stats": {
            "files": len(files),
            "documents": total_docs,
            "images": total_images,
            "media": total_media,
            "text_chars": total_text_chars,
            "sections": len(sections),
        },
    }
    log(
        f"资料读取完成：{len(sections)} 个分组，"
        f"{total_docs} 个文档，{total_images} 张图片，{total_media} 个音视频。"
    )
    return inventory


def shape_rgb(shape: Any) -> tuple[int, int, int] | None:
    try:
        fill = shape.fill
        if not fill:
            return None
        rgb = fill.fore_color.rgb
        return rgb_obj_to_tuple(rgb)
    except Exception:
        return None


def run_rgb(run: Any) -> tuple[int, int, int] | None:
    try:
        return rgb_obj_to_tuple(run.font.color.rgb)
    except Exception:
        return None


def iter_pptx_text_runs(shape: Any) -> list[Any]:
    runs: list[Any] = []
    if not getattr(shape, "has_text_frame", False):
        return runs
    try:
        for para in shape.text_frame.paragraphs:
            runs.extend(list(para.runs))
    except Exception:
        return runs
    return runs


def font_name_from_run_xml(run: Any, tag: str) -> str | None:
    try:
        r_pr = run._r.rPr
        if r_pr is None:
            return None
        node = r_pr.find(qn(f"a:{tag}"))
        if node is None:
            return None
        value = node.get("typeface")
        if not value or value.startswith("+"):
            return None
        return value
    except Exception:
        return None


def run_font_names(run: Any) -> list[tuple[str, str]]:
    names: list[tuple[str, str]] = []
    east_asian = font_name_from_run_xml(run, "ea")
    latin = getattr(run.font, "name", None) or font_name_from_run_xml(run, "latin")
    complex_script = font_name_from_run_xml(run, "cs")
    if east_asian:
        names.append(("east_asian", east_asian))
    if latin and not str(latin).startswith("+"):
        names.append(("latin", str(latin)))
    if complex_script:
        names.append(("complex", complex_script))
    return names


def set_run_font(run: Any, font_name: str) -> None:
    if not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r_pr = run._r.get_or_add_rPr()
        for tag in ("latin", "ea", "cs"):
            node = r_pr.find(qn(f"a:{tag}"))
            if node is None:
                node = OxmlElement(f"a:{tag}")
                r_pr.append(node)
            node.set("typeface", font_name)
    except Exception:
        pass


def median(values: list[float], fallback: float) -> float:
    if not values:
        return fallback
    values = sorted(values)
    mid = len(values) // 2
    if len(values) % 2:
        return values[mid]
    return (values[mid - 1] + values[mid]) / 2


def default_profile() -> dict[str, Any]:
    return {
        "version": 1,
        "source_template": "",
        "template_type": "none",
        "slide_width": 12192000,
        "slide_height": 6858000,
        "slide_width_in": 13.333,
        "slide_height_in": 7.5,
        "font": {
            "primary": DEFAULT_FONT,
            "title_size": 34,
            "subtitle_size": 17,
            "body_size": 16,
            "caption_size": 10,
        },
        "palette": {
            "background": "#F7F8FB",
            "surface": "#FFFFFF",
            "primary": "#1F2937",
            "accent": "#2F80ED",
            "accent_2": "#16A085",
            "muted": "#6B7280",
        },
        "layout": {
            "title_box": [0.72, 0.48, 11.9, 0.72],
            "body_box": [0.78, 1.38, 11.78, 5.22],
            "picture_box": [7.45, 1.56, 4.84, 4.54],
            "media_box": [0.78, 1.38, 11.78, 5.22],
            "left_text_box": [0.82, 1.48, 5.85, 4.86],
        },
        "modules": {
            "font": True,
            "color": True,
            "background": True,
            "logo": False,
            "layout": True,
            "page_size": True,
            "picture_style": True,
        },
        "assets": {"logos": [], "background_image": None},
        "analysis": {
            "slide_count": 0,
            "layout_count": 0,
            "detected_fonts": [],
            "detected_colors": [],
            "detected_logo_count": 0,
            "notes": [],
        },
        "slide_blueprints": [],
    }


def choose_palette(colors: Counter[tuple[int, int, int]]) -> dict[str, str]:
    if not colors:
        return default_profile()["palette"]

    ranked = [color for color, _ in colors.most_common(80)]
    dark = [c for c in ranked if color_brightness(c) < 0.35]
    light = [c for c in ranked if color_brightness(c) > 0.82]
    accents = [c for c in ranked if 0.22 < color_brightness(c) < 0.78 and color_saturation(c) > 0.12]

    primary = dark[0] if dark else (31, 41, 55)
    background = light[0] if light else (247, 248, 251)
    accent = accents[0] if accents else (47, 128, 237)
    accent_2 = accents[1] if len(accents) > 1 else (22, 160, 133)
    surface = (255, 255, 255) if color_brightness(background) < 0.72 else lighten(background, 0.7)
    muted = lighten(primary, 0.35)
    return {
        "background": tuple_to_hex(background),
        "surface": tuple_to_hex(surface),
        "primary": tuple_to_hex(primary),
        "accent": tuple_to_hex(accent),
        "accent_2": tuple_to_hex(accent_2),
        "muted": tuple_to_hex(muted),
    }


def fit_rect_from_samples(samples: list[list[float]], fallback: list[float]) -> list[float]:
    if not samples:
        return fallback
    return [
        round(median([item[0] for item in samples], fallback[0]), 3),
        round(median([item[1] for item in samples], fallback[1]), 3),
        round(median([item[2] for item in samples], fallback[2]), 3),
        round(median([item[3] for item in samples], fallback[3]), 3),
    ]


def shape_rect_inches(shape: Any) -> list[float]:
    return [
        round(emu_to_inches(shape.left), 3),
        round(emu_to_inches(shape.top), 3),
        round(emu_to_inches(shape.width), 3),
        round(emu_to_inches(shape.height), 3),
    ]


def is_full_slide_shape(shape: Any, slide_w: int, slide_h: int) -> bool:
    try:
        area_ratio = (int(shape.width) * int(shape.height)) / max(1, slide_w * slide_h)
        covers_width = int(shape.width) >= slide_w * 0.92
        covers_height = int(shape.height) >= slide_h * 0.92
        near_origin = int(shape.left) <= slide_w * 0.04 and int(shape.top) <= slide_h * 0.04
        return area_ratio > 0.78 and covers_width and covers_height and near_origin
    except Exception:
        return False


def max_text_size(shape: Any) -> float:
    sizes: list[float] = []
    for run in iter_pptx_text_runs(shape):
        size = getattr(run.font, "size", None)
        if size:
            sizes.append(float(size.pt))
    return max(sizes) if sizes else 0


def slide_blueprint(slide: Any, slide_index: int, slide_w: int, slide_h: int) -> dict[str, Any]:
    text_slots: list[dict[str, Any]] = []
    picture_slots: list[dict[str, Any]] = []
    for shape_index, shape in enumerate(slide.shapes):
        rect = shape_rect_inches(shape)
        if getattr(shape, "has_text_frame", False):
            text = getattr(shape, "text", "").strip()
            if text and not is_full_slide_shape(shape, slide_w, slide_h):
                size = max_text_size(shape)
                role = "body"
                if shape.top < slide_h * 0.24 and (size >= 24 or shape.height < slide_h * 0.18):
                    role = "title"
                elif shape.top > slide_h * 0.84:
                    role = "footer"
                text_slots.append(
                    {
                        "index": shape_index,
                        "role": role,
                        "rect": rect,
                        "text_preview": text[:120],
                        "font_size": round(size, 1) if size else None,
                    }
                )
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            picture_slots.append(
                {
                    "index": shape_index,
                    "rect": rect,
                    "is_background": is_full_slide_shape(shape, slide_w, slide_h),
                }
            )
    return {"index": slide_index + 1, "text_slots": text_slots, "picture_slots": picture_slots}


def extract_logo_assets(shape: Any, slide_w: int, slide_h: int, asset_dir: Path, index: int) -> dict[str, Any] | None:
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return None
        area_ratio = (shape.width * shape.height) / max(1, slide_w * slide_h)
        near_edge = (
            shape.left < slide_w * 0.15
            or shape.left + shape.width > slide_w * 0.85
            or shape.top < slide_h * 0.15
            or shape.top + shape.height > slide_h * 0.85
        )
        if area_ratio > 0.07 or not near_edge:
            return None
        image = shape.image
        ext = image.ext or "png"
        asset_path = asset_dir / f"logo_{index}.{ext}"
        asset_path.write_bytes(image.blob)
        return {
            "path": str(asset_path),
            "x": emu_to_inches(shape.left),
            "y": emu_to_inches(shape.top),
            "w": emu_to_inches(shape.width),
            "h": emu_to_inches(shape.height),
        }
    except Exception:
        return None


def learn_pptx_template(template_path: Path, output_dir: Path, log: LogFn) -> dict[str, Any]:
    profile = default_profile()
    profile["source_template"] = str(template_path)
    profile["template_type"] = "pptx"

    prs = Presentation(str(template_path))
    profile["slide_width"] = int(prs.slide_width)
    profile["slide_height"] = int(prs.slide_height)
    profile["slide_width_in"] = round(emu_to_inches(prs.slide_width), 3)
    profile["slide_height_in"] = round(emu_to_inches(prs.slide_height), 3)

    asset_dir = ensure_dir(output_dir / "_template_assets" / safe_filename(template_path.stem))
    fonts: Counter[str] = Counter()
    east_asian_fonts: Counter[str] = Counter()
    font_sizes: list[float] = []
    title_sizes: list[float] = []
    colors: Counter[tuple[int, int, int]] = Counter()
    title_boxes: list[list[float]] = []
    body_boxes: list[list[float]] = []
    picture_boxes: list[list[float]] = []
    media_boxes: list[list[float]] = []
    logos: list[dict[str, Any]] = []
    logo_seen: set[tuple[int, int, int, int]] = set()
    blueprints: list[dict[str, Any]] = []

    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    slide_count = len(prs.slides)

    for slide_index, slide in enumerate(prs.slides):
        if slide_index >= 30:
            break
        blueprints.append(slide_blueprint(slide, slide_index, slide_w, slide_h))
        for shape_index, shape in enumerate(slide.shapes):
            fill_rgb = shape_rgb(shape)
            if fill_rgb:
                colors[fill_rgb] += 2

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_rect = shape_rect_inches(shape)
                if not is_full_slide_shape(shape, slide_w, slide_h):
                    picture_boxes.append(picture_rect)
                    if (shape.width * shape.height) / max(1, slide_w * slide_h) > 0.08:
                        media_boxes.append(picture_rect)
                logo_key = (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
                if logo_key not in logo_seen and len(logos) < 4:
                    logo = extract_logo_assets(shape, slide_w, slide_h, asset_dir, len(logos) + 1)
                    if logo:
                        logos.append(logo)
                        logo_seen.add(logo_key)

            if not getattr(shape, "has_text_frame", False):
                continue
            text = getattr(shape, "text", "").strip()
            if not text:
                continue
            shape_box = [
                emu_to_inches(shape.left),
                emu_to_inches(shape.top),
                emu_to_inches(shape.width),
                emu_to_inches(shape.height),
            ]
            shape_sizes: list[float] = []
            for run in iter_pptx_text_runs(shape):
                for kind, name in run_font_names(run):
                    if kind == "east_asian":
                        east_asian_fonts[name] += 3
                    fonts[name] += 1
                size = getattr(run.font, "size", None)
                if size:
                    size_pt = size.pt
                    font_sizes.append(size_pt)
                    shape_sizes.append(size_pt)
                rgb = run_rgb(run)
                if rgb:
                    colors[rgb] += 3

            local_size = max(shape_sizes) if shape_sizes else 0
            is_top = shape.top < slide_h * 0.28
            is_big = local_size >= 24 or (shape.height < slide_h * 0.2 and len(text) <= 80)
            if is_top and is_big:
                title_boxes.append(shape_box)
                if local_size:
                    title_sizes.append(local_size)
            else:
                body_boxes.append(shape_box)
                if shape.height > slide_h * 0.26 and shape.width > slide_w * 0.36:
                    media_boxes.append(shape_box)

    primary_font = east_asian_fonts.most_common(1)[0][0] if east_asian_fonts else (fonts.most_common(1)[0][0] if fonts else DEFAULT_FONT)
    title_size = int(round(median(title_sizes, 34)))
    body_candidates = [size for size in font_sizes if size < max(title_size, 24)]
    body_size = int(round(median(body_candidates, 16)))

    fallback_layout = default_profile()["layout"]
    profile["font"] = {
        "primary": primary_font,
        "title_size": int(clamp(title_size, 24, 46)),
        "subtitle_size": int(clamp(body_size + 2, 13, 22)),
        "body_size": int(clamp(body_size, 11, 21)),
        "caption_size": int(clamp(body_size - 4, 8, 12)),
    }
    profile["palette"] = choose_palette(colors)
    profile["layout"] = {
        "title_box": fit_rect_from_samples(title_boxes, fallback_layout["title_box"]),
        "body_box": fit_rect_from_samples(body_boxes, fallback_layout["body_box"]),
        "picture_box": fit_rect_from_samples(picture_boxes, fallback_layout["picture_box"]),
        "media_box": fit_rect_from_samples(media_boxes, fallback_layout["media_box"]),
        "left_text_box": fallback_layout["left_text_box"],
    }
    profile["assets"]["logos"] = logos
    profile["modules"]["logo"] = bool(logos)
    profile["analysis"] = {
        "slide_count": slide_count,
        "layout_count": len(prs.slide_layouts),
        "detected_fonts": [name for name, _ in fonts.most_common(8)],
        "detected_colors": [tuple_to_hex(color) for color, _ in colors.most_common(12)],
        "detected_logo_count": len(logos),
        "notes": [
            "PPTX 模板已解析：尺寸、中文字体、字号、颜色、逐页文本槽位、图片/媒体展示区和角落 Logo/水印。",
            "生成 PPTX 时会优先复用模板页面骨架，并按模板槽位填充文字、图片和音视频。",
        ],
    }
    profile["slide_blueprints"] = blueprints
    return profile


def learn_pdf_template(template_path: Path, output_dir: Path, log: LogFn) -> dict[str, Any]:
    profile = default_profile()
    profile["source_template"] = str(template_path)
    profile["template_type"] = "pdf"
    asset_dir = ensure_dir(output_dir / "_template_assets" / safe_filename(template_path.stem))

    with fitz.open(template_path) as pdf:
        if len(pdf) == 0:
            raise ValueError("PDF 模板没有页面。")
        page = pdf[0]
        rect = page.rect
        profile["slide_width"] = points_to_emu(rect.width)
        profile["slide_height"] = points_to_emu(rect.height)
        profile["slide_width_in"] = round(rect.width / POINTS_PER_INCH, 3)
        profile["slide_height_in"] = round(rect.height / POINTS_PER_INCH, 3)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        bg_path = asset_dir / "pdf_template_background.png"
        pix.save(bg_path)
        text = page.get_text()

    lines = split_meaningful_lines(text)
    if lines:
        profile["font"]["title_size"] = 32
        profile["font"]["body_size"] = 15
    profile["assets"]["background_image"] = str(bg_path)
    profile["analysis"] = {
        "slide_count": 1,
        "layout_count": 1,
        "detected_fonts": [DEFAULT_FONT],
        "detected_colors": [],
        "detected_logo_count": 0,
        "notes": [
            "PDF 模板以首屏背景图方式保留底图和视觉氛围。",
            "PDF 原文语义不如 PPTX 完整，建议优先选择 PPTX 模板以获得更可编辑的结果。",
        ],
    }
    return profile


def learn_template(template_path: Path, output_dir: Path, log: LogFn = noop_log) -> dict[str, Any]:
    template_path = template_path.resolve()
    output_dir = ensure_dir(output_dir.resolve())
    if template_path.suffix.lower() not in SUPPORTED_TEMPLATE_EXTS:
        raise ValueError("模板仅支持 .pptx。")
    if not template_path.exists():
        raise FileNotFoundError(f"模板不存在：{template_path}")

    log(f"学习模板：{template_path}")
    profile = learn_pptx_template(template_path, output_dir, log)

    profile_path = output_dir / f"{safe_filename(template_path.stem)}_template_profile.json"
    profile_path.write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
    log(f"模板学习完成：{profile_path}")
    return profile


def compact_sections_for_model(sections: list[dict[str, Any]], limit: int = 18) -> list[dict[str, Any]]:
    compact: list[dict[str, Any]] = []
    for section in sections[:limit]:
        compact.append(
            {
                "title": section["title"],
                "bullets": section["bullets"][:6],
                "image_count": len(section["images"]),
                "media_count": len(section["media"]),
                "documents": [doc["name"] for doc in section["documents"][:5]],
            }
        )
    return compact


def parse_json_object(text: str) -> Any:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?", "", text).strip()
        text = re.sub(r"```$", "", text).strip()
    try:
        return json.loads(text)
    except Exception:
        pass
    match = re.search(r"(\{.*\}|\[.*\])", text, flags=re.S)
    if not match:
        raise ValueError("模型返回内容不是 JSON。")
    return json.loads(match.group(1))


def call_model(prompt: dict[str, Any], config: dict[str, str], log: LogFn, purpose: str) -> str:
    mode = config.get("mode", "规则引擎")
    model_name = config.get("model", "").strip() or ("qwen2.5" if mode == "Ollama 本地模型" else "gpt-4.1-mini")
    base_url = config.get("base_url", "").strip()
    api_key = config.get("api_key", "").strip()
    if not base_url:
        base_url = "http://localhost:11434/api/chat" if mode == "Ollama 本地模型" else "https://api.openai.com/v1/chat/completions"

    log(f"调用模型：{purpose} / {mode} / {model_name}")
    if mode == "Ollama 本地模型":
        response = requests.post(
            base_url,
            json={
                "model": model_name,
                "messages": [{"role": "user", "content": json.dumps(prompt, ensure_ascii=False)}],
                "stream": False,
                "options": {"temperature": 0.2},
            },
            timeout=80,
        )
        response.raise_for_status()
        data = response.json()
        return data.get("message", {}).get("content", "")

    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    response = requests.post(
        base_url,
        headers=headers,
        json={
            "model": model_name,
            "messages": [{"role": "user", "content": json.dumps(prompt, ensure_ascii=False)}],
            "temperature": 0.2,
        },
        timeout=80,
    )
    response.raise_for_status()
    data = response.json()
    return data["choices"][0]["message"]["content"]


def refine_template_with_model(profile: dict[str, Any], config: dict[str, str], log: LogFn) -> dict[str, Any]:
    mode = config.get("mode", "规则引擎")
    if mode == "规则引擎":
        return profile

    prompt = {
        "task": "请根据这个PPTX模板解析结果，推断模板视觉风格和生成建议。只返回JSON，不要解释。",
        "schema": {
            "style_summary": "一句话总结模板风格",
            "recommended_keep_modules": ["font", "color", "background", "logo", "layout", "page_size", "picture_style"],
            "design_notes": ["给生成器的短建议，每条不超过30字，最多5条"],
        },
        "template_profile": {
            "template_type": profile.get("template_type"),
            "slide_size": [profile.get("slide_width_in"), profile.get("slide_height_in")],
            "font": profile.get("font"),
            "palette": profile.get("palette"),
            "layout": profile.get("layout"),
            "analysis": profile.get("analysis"),
        },
    }
    try:
        content = call_model(prompt, config, log, "理解模板风格")
        refined = parse_json_object(content)
        if isinstance(refined, dict):
            analysis = profile.setdefault("analysis", {})
            if refined.get("style_summary"):
                analysis["model_style_summary"] = str(refined["style_summary"])[:160]
            if refined.get("recommended_keep_modules"):
                allowed = set(default_profile()["modules"].keys())
                analysis["model_recommended_keep_modules"] = [
                    str(item) for item in refined["recommended_keep_modules"] if str(item) in allowed
                ]
            if refined.get("design_notes"):
                analysis["model_design_notes"] = [str(item)[:80] for item in refined["design_notes"][:5]]
        log("模型模板理解完成，已写入模板分析结果。")
    except Exception as exc:
        profile.setdefault("analysis", {}).setdefault("notes", []).append(f"模型模板理解失败，已使用本地解析：{exc}")
        log(f"模型模板理解失败，继续使用本地解析：{exc}")
    return profile


def refine_with_model(inventory: dict[str, Any], config: dict[str, str], log: LogFn) -> dict[str, Any]:
    mode = config.get("mode", "规则引擎")
    if mode == "规则引擎":
        return inventory

    prompt = {
        "task": "请把这些周报资料整理成更适合PPT展示的中文短标题和要点。只返回JSON，不要解释。",
        "schema": {
            "sections": [
                {
                    "title": "模块名，不超过18字",
                    "bullets": ["每条不超过32字，适合PPT展示，最多6条"],
                }
            ]
        },
        "source": compact_sections_for_model(inventory["sections"]),
    }
    try:
        content = call_model(prompt, config, log, "整理资料内容")
        refined = parse_json_object(content)
        model_sections = refined["sections"] if isinstance(refined, dict) else refined
        for source_section, model_section in zip(inventory["sections"], model_sections):
            if isinstance(model_section, dict):
                title = str(model_section.get("title") or "").strip()
                bullets = [clean_line(str(item)) for item in model_section.get("bullets", []) if clean_line(str(item))]
                if title:
                    source_section["title"] = title[:28]
                if bullets:
                    source_section["bullets"] = bullets[:7]
        log("模型整理完成，已应用到页面计划。")
    except Exception as exc:
        inventory["warnings"].append(f"模型调用失败，已回退到规则引擎：{exc}")
        log(f"模型调用失败，继续使用规则引擎：{exc}")
    return inventory


def chunk_list(items: list[Any], size: int) -> list[list[Any]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def build_slide_plan(
    inventory: dict[str, Any],
    profile: dict[str, Any],
    title: str,
    optimize: bool,
    preserve: set[str],
) -> list[dict[str, Any]]:
    stats = inventory["stats"]
    sections = inventory["sections"]
    slides: list[dict[str, Any]] = []
    today = dt.datetime.now().strftime("%Y-%m-%d")
    subtitle = (
        f"{stats['documents']} 个文档 / {stats['images']} 张图片 / "
        f"{stats['media']} 个音视频 / {stats['sections']} 个分组"
    )

    slides.append(
        {
            "kind": "cover",
            "title": title,
            "subtitle": subtitle,
            "date": today,
            "optimize": optimize,
        }
    )

    contents = [section["title"] for section in sections[:12]]
    slides.append(
        {
            "kind": "agenda",
            "title": "内容结构",
            "bullets": contents or ["暂无可解析分组"],
        }
    )

    summary_bullets = [
        f"共整理 {stats['sections']} 个资料分组，生成结构化页面。",
        f"解析 {stats['documents']} 个文档，提取约 {stats['text_chars']} 个字符。",
        f"收集 {stats['images']} 张图片和 {stats['media']} 个音视频素材。",
    ]
    if inventory.get("warnings"):
        summary_bullets.append(f"有 {len(inventory['warnings'])} 条解析提醒，详见日志。")
    slides.append({"kind": "summary", "title": "资料概览", "bullets": summary_bullets})

    for section in sections:
        bullets = section["bullets"][:]
        if not bullets:
            bullets = ["暂无可解析正文，已保留素材清单。"]
        image_chunks = chunk_list(section["images"], 2) or [[]]
        bullet_chunks = chunk_list(bullets, 6)
        max_chunks = max(len(image_chunks), len(bullet_chunks))
        for index in range(max_chunks):
            slide_bullets = bullet_chunks[index] if index < len(bullet_chunks) else []
            slide_images = image_chunks[index] if index < len(image_chunks) else []
            title_suffix = f" {index + 1}" if max_chunks > 1 else ""
            slides.append(
                {
                    "kind": "section",
                    "title": section["title"] + title_suffix,
                    "bullets": slide_bullets,
                    "images": slide_images,
                    "documents": [doc["name"] for doc in section["documents"][:4]],
                    "media": [],
                }
            )
        if section["media"]:
            for media_item in section["media"]:
                media_kind = media_kind_from_path(Path(media_item["path"]))
                media_title = "音频展示" if media_kind == "audio" else "产出展示"
                slides.append(
                    {
                        "kind": "media",
                        "title": f"{section['title']} {media_title}",
                        "bullets": [],
                        "images": [],
                        "documents": [doc["name"] for doc in section["documents"][:4]],
                        "media": [media_item],
                    }
                )

    media_items: list[str] = []
    for section in sections:
        for item in section["media"]:
            media_items.append(f"{section['title']} / {item['name']}")
    for index, chunk in enumerate(chunk_list(media_items, 10)):
        slides.append({"kind": "text", "title": "音视频素材清单" if index == 0 else "音视频素材清单续", "bullets": chunk})

    slides.append(
        {
            "kind": "ending",
            "title": "下一步",
            "bullets": [
                "人工检查重点页的标题和配图。",
                "按需要替换无法自动解析的旧版 .doc 或特殊媒体。",
                "确认后可在 PowerPoint 中按需另存为 PDF。",
            ],
        }
    )
    return slides


def get_preserve_set(options: dict[str, Any]) -> set[str]:
    return {key for key, value in options.get("preserve", {}).items() if value}


def remove_all_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[0]


def find_blank_layout(prs: Presentation) -> Any:
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[-1]


def create_presentation(template_path: Path | None, profile: dict[str, Any], preserve: set[str]) -> Presentation:
    if template_path and template_path.suffix.lower() == ".pptx" and "layout" in preserve:
        prs = Presentation(str(template_path))
        remove_all_slides(prs)
    else:
        prs = Presentation()
    if "page_size" in preserve or not template_path:
        prs.slide_width = int(profile["slide_width"])
        prs.slide_height = int(profile["slide_height"])
    return prs


def rect_to_emu(rect: list[float]) -> tuple[Emu, Emu, Emu, Emu]:
    return tuple(Emu(inches_to_emu(float(item))) for item in rect)  # type: ignore[return-value]


def box(x: float, y: float, w: float, h: float) -> list[float]:
    return [x, y, w, h]


def add_rect(slide: Any, rect: list[float], color: tuple[int, int, int], transparency: int = 0, line: bool = False) -> Any:
    left, top, width, height = rect_to_emu(rect)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = tuple_to_rgb(color)
    if transparency:
        shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = tuple_to_rgb(darken(color, 0.16))
    else:
        shape.line.fill.background()
    return shape


def set_slide_background(slide: Any, color: tuple[int, int, int]) -> None:
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = tuple_to_rgb(color)
    except Exception:
        add_rect(slide, [0, 0, 50, 50], color)


def add_textbox(
    slide: Any,
    text: str,
    rect: list[float],
    font_name: str,
    font_size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    align: str = "left",
) -> Any:
    left, top, width, height = rect_to_emu(rect)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    if align == "center":
        paragraph.alignment = 2
    elif align == "right":
        paragraph.alignment = 3
    for run in paragraph.runs:
        set_run_font(run, font_name)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = tuple_to_rgb(color)
    return shape


def add_bullets(
    slide: Any,
    bullets: list[str],
    rect: list[float],
    font_name: str,
    font_size: int,
    color: tuple[int, int, int],
    accent: tuple[int, int, int],
) -> Any:
    left, top, width, height = rect_to_emu(rect)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    if not bullets:
        bullets = ["暂无正文内容。"]
    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.space_after = Pt(7)
        paragraph.level = 0
        for run in paragraph.runs:
            set_run_font(run, font_name)
            run.font.size = Pt(font_size)
            run.font.color.rgb = tuple_to_rgb(color)
    return shape


def add_label(
    slide: Any,
    text: str,
    rect: list[float],
    font_name: str,
    color: tuple[int, int, int],
    bg: tuple[int, int, int],
) -> None:
    add_rect(slide, rect, bg, transparency=0)
    add_textbox(slide, text, [rect[0] + 0.08, rect[1] + 0.03, rect[2] - 0.16, rect[3] - 0.02], font_name, 9, color, False, "center")


def image_size(path: Path) -> tuple[int, int] | None:
    try:
        with Image.open(path) as img:
            return img.size
    except Exception:
        return None


def add_picture_contained(slide: Any, path: Path, rect: list[float]) -> bool:
    size = image_size(path)
    if not size:
        return False
    img_w, img_h = size
    if img_w <= 0 or img_h <= 0:
        return False
    box_w = inches_to_emu(rect[2])
    box_h = inches_to_emu(rect[3])
    scale = min(box_w / img_w, box_h / img_h)
    width = int(img_w * scale)
    height = int(img_h * scale)
    left = inches_to_emu(rect[0]) + int((box_w - width) / 2)
    top = inches_to_emu(rect[1]) + int((box_h - height) / 2)
    try:
        slide.shapes.add_picture(str(path), Emu(left), Emu(top), width=Emu(width), height=Emu(height))
        return True
    except Exception:
        return False


def media_kind_from_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in VIDEO_EXTS:
        return "video"
    if suffix in AUDIO_EXTS:
        return "audio"
    return "media"


def normalize_media_items(items: list[Any]) -> list[dict[str, str]]:
    normalized: list[dict[str, str]] = []
    for item in items or []:
        if isinstance(item, dict):
            path = str(item.get("path") or "")
            name = str(item.get("name") or (Path(path).name if path else "媒体素材"))
            normalized.append({"path": path, "name": name, "relative": str(item.get("relative") or name)})
        else:
            name = str(item)
            normalized.append({"path": "", "name": name, "relative": name})
    return normalized


def grid_rects(rect: list[float], count: int, prefer_video: bool = False) -> list[list[float]]:
    count = max(1, min(count, 6))
    x, y, w, h = rect
    gap = 0.16
    if count == 1:
        return [[x, y, w, h]]
    if prefer_video and count <= 2:
        rows, cols = count, 1
    elif count <= 4:
        rows, cols = 2, 2
    else:
        rows, cols = 2, 3
    cell_w = (w - gap * (cols - 1)) / cols
    cell_h = (h - gap * (rows - 1)) / rows
    rects: list[list[float]] = []
    for row in range(rows):
        for col in range(cols):
            if len(rects) >= count:
                return rects
            rects.append([x + col * (cell_w + gap), y + row * (cell_h + gap), cell_w, cell_h])
    return rects


def add_media_card(slide: Any, media: dict[str, str], rect: list[float], profile: dict[str, Any]) -> None:
    font = profile["font"].get("primary") or DEFAULT_FONT
    palette = profile.get("palette", {})
    primary = hex_to_tuple(palette.get("primary"), (31, 41, 55))
    surface = hex_to_tuple(palette.get("surface"), (255, 255, 255))
    accent = hex_to_tuple(palette.get("accent"), (47, 128, 237))
    kind = media_kind_from_path(Path(media.get("path") or media.get("name") or ""))
    fill = darken(surface, 0.02) if kind == "audio" else (34, 39, 49)
    line = lighten(primary, 0.55) if kind == "audio" else darken(accent, 0.18)
    card = add_rect(slide, rect, fill, transparency=0, line=True)
    try:
        card.line.color.rgb = tuple_to_rgb(line)
    except Exception:
        pass
    media_path = media.get("path") or ""
    if media_path:
        try:
            card.click_action.hyperlink.address = media_path
        except Exception:
            pass
    symbol = "▶" if kind == "video" else "♪"
    symbol_color = readable_text_color(fill)
    if kind == "audio":
        symbol_color = accent
    symbol_box = add_textbox(slide, symbol, [rect[0] + 0.15, rect[1] + 0.12, 0.46, 0.34], font, 16, symbol_color, True, "center")
    label_box = add_textbox(
        slide,
        media.get("name", "媒体素材"),
        [rect[0] + 0.68, rect[1] + 0.13, max(0.5, rect[2] - 0.86), max(0.28, rect[3] - 0.22)],
        font,
        11 if rect[3] > 0.75 else 9,
        symbol_color if kind == "video" else primary,
        False,
        "left",
    )
    if media_path:
        for shape in (symbol_box, label_box):
            try:
                shape.click_action.hyperlink.address = media_path
            except Exception:
                pass


def set_shape_hyperlink(shape: Any, address: str) -> None:
    if not address:
        return
    try:
        shape.click_action.hyperlink.address = address
    except Exception:
        pass


def media_box(profile: dict[str, Any]) -> list[float]:
    rect = profile["layout"].get("media_box") or profile["layout"].get("body_box") or default_profile()["layout"]["media_box"]
    width = float(profile["slide_width_in"])
    height = float(profile["slide_height_in"])
    x = clamp(float(rect[0]), 0.2, max(0.2, width - 1.0))
    y = clamp(float(rect[1]), 0.7, max(0.7, height - 1.0))
    w = clamp(float(rect[2]), 2.0, max(2.0, width - x - 0.35))
    h = clamp(float(rect[3]), 1.0, max(1.0, height - y - 0.45))
    return [x, y, w, h]


def plan_media_on_slide(slide: Any, item: dict[str, Any], profile: dict[str, Any], body_shape: Any | None, slide_number: int) -> list[dict[str, Any]]:
    media_items = [item for item in normalize_media_items(item.get("media", [])) if item.get("path")]
    if not media_items:
        return []
    rect = shape_rect_inches(body_shape) if body_shape is not None else media_box(profile)
    rect[2] = max(1.35, rect[2])
    rect[3] = max(0.8, rect[3])
    insertions: list[dict[str, Any]] = []
    for media in media_items[:1]:
        if body_shape is not None:
            set_shape_hyperlink(body_shape, media["path"])
        insertions.append(
            {
                "slide": slide_number,
                "path": media["path"],
                "name": media["name"],
                "kind": media_kind_from_path(Path(media["path"])),
                "rect": rect,
                "mode": "single_link_text",
            }
        )
    return insertions


def draw_template_logos(slide: Any, profile: dict[str, Any], preserve: set[str]) -> None:
    if "logo" not in preserve:
        return
    for logo in profile.get("assets", {}).get("logos", [])[:3]:
        path = Path(logo["path"])
        if not path.exists():
            continue
        rect = [float(logo["x"]), float(logo["y"]), float(logo["w"]), float(logo["h"])]
        add_picture_contained(slide, path, rect)


def apply_slide_design(slide: Any, profile: dict[str, Any], preserve: set[str], optimize: bool) -> dict[str, tuple[int, int, int]]:
    palette = profile["palette"]
    background = hex_to_tuple(palette.get("background"), (247, 248, 251))
    surface = hex_to_tuple(palette.get("surface"), (255, 255, 255))
    primary = hex_to_tuple(palette.get("primary"), (31, 41, 55))
    accent = hex_to_tuple(palette.get("accent"), (47, 128, 237))
    accent_2 = hex_to_tuple(palette.get("accent_2"), (22, 160, 133))
    muted = hex_to_tuple(palette.get("muted"), (107, 114, 128))

    if "background" in preserve:
        bg_image = profile.get("assets", {}).get("background_image")
        if bg_image and Path(bg_image).exists():
            add_picture_contained(
                slide,
                Path(bg_image),
                [0, 0, float(profile["slide_width_in"]), float(profile["slide_height_in"])],
            )
        else:
            set_slide_background(slide, background)
    else:
        background = (245, 247, 250)
        surface = (255, 255, 255)
        set_slide_background(slide, background)

    if optimize:
        slide_w = float(profile["slide_width_in"])
        slide_h = float(profile["slide_height_in"])
        if "background" not in preserve:
            set_slide_background(slide, (244, 247, 249))
            add_rect(slide, [0, 0, slide_w, 0.22], accent)
            add_rect(slide, [0, slide_h - 0.18, slide_w, 0.18], accent_2)
        add_rect(slide, [0.52, 0.52, 0.08, 5.95], accent, transparency=0)
    else:
        add_rect(slide, [0.65, 1.22, 1.35, 0.05], accent, transparency=0)

    draw_template_logos(slide, profile, preserve)
    return {
        "background": background,
        "surface": surface,
        "primary": primary,
        "accent": accent,
        "accent_2": accent_2,
        "muted": muted,
    }


def title_box(profile: dict[str, Any]) -> list[float]:
    rect = profile["layout"].get("title_box") or default_profile()["layout"]["title_box"]
    width = float(profile["slide_width_in"])
    return [clamp(rect[0], 0.45, 1.2), clamp(rect[1], 0.25, 0.95), clamp(rect[2], 5, width - 1.1), clamp(rect[3], 0.45, 1.05)]


def body_box(profile: dict[str, Any]) -> list[float]:
    width = float(profile["slide_width_in"])
    height = float(profile["slide_height_in"])
    rect = profile["layout"].get("body_box") or default_profile()["layout"]["body_box"]
    return [clamp(rect[0], 0.7, 1.3), clamp(rect[1], 1.35, 2.1), clamp(rect[2], 5.5, width - 1.5), clamp(rect[3], 3.8, height - 2.0)]


def render_cover(slide: Any, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], preserve: set[str]) -> None:
    font = profile["font"]["primary"] if "font" in preserve else DEFAULT_FONT
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])
    title_size = profile["font"]["title_size"] if "font" in preserve else 38
    primary = colors["primary"]
    accent = colors["accent"]
    surface = colors["surface"]

    add_rect(slide, [0.72, 1.28, 0.12, 2.28], accent)
    add_textbox(slide, item["title"], [1.05, 1.18, slide_w - 2.0, 1.45], font, title_size, primary, True)
    add_textbox(slide, item["subtitle"], [1.08, 2.75, slide_w - 2.0, 0.55], font, 15, colors["muted"])
    add_textbox(slide, item["date"], [1.08, 3.36, 4.5, 0.45], font, 12, colors["muted"])
    add_rect(slide, [1.08, slide_h - 1.45, slide_w - 2.16, 0.02], colors["muted"], transparency=45)
    add_textbox(slide, "由模板学习 + 资料整理自动生成", [1.08, slide_h - 1.18, 5.5, 0.35], font, 10, colors["muted"])
    if item.get("optimize"):
        add_label(slide, "优化设计版", [slide_w - 2.15, slide_h - 1.22, 1.35, 0.34], font, readable_text_color(accent), accent)
    else:
        add_label(slide, "同风格版", [slide_w - 2.0, slide_h - 1.22, 1.2, 0.34], font, primary, surface)


def render_text_slide(slide: Any, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], preserve: set[str], optimize: bool) -> None:
    font = profile["font"]["primary"] if "font" in preserve else DEFAULT_FONT
    title_size = profile["font"]["title_size"] - 5 if "font" in preserve else 28
    body_size = profile["font"]["body_size"] if "font" in preserve else 15
    tbox = title_box(profile)
    bbox = body_box(profile)
    add_textbox(slide, item["title"], tbox, font, int(clamp(title_size, 22, 36)), colors["primary"], True)
    if optimize:
        add_rect(slide, [bbox[0] - 0.1, bbox[1] - 0.12, bbox[2] + 0.2, bbox[3] + 0.18], colors["surface"], transparency=0, line=True)
        bbox = [bbox[0] + 0.18, bbox[1] + 0.18, bbox[2] - 0.36, bbox[3] - 0.28]
    add_bullets(slide, item.get("bullets", []), bbox, font, int(clamp(body_size, 12, 19)), colors["primary"], colors["accent"])


def render_section_slide(slide: Any, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], preserve: set[str], optimize: bool) -> None:
    font = profile["font"]["primary"] if "font" in preserve else DEFAULT_FONT
    title_size = profile["font"]["title_size"] - 6 if "font" in preserve else 27
    body_size = profile["font"]["body_size"] if "font" in preserve else 15
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])

    add_textbox(slide, item["title"], title_box(profile), font, int(clamp(title_size, 21, 34)), colors["primary"], True)
    images = item.get("images", [])[:2]
    has_images = bool(images)

    if has_images:
        text_rect = [0.88, 1.55, slide_w * 0.46, slide_h - 2.1]
        image_area = [slide_w * 0.56, 1.42, slide_w * 0.38, slide_h - 2.05]
        add_bullets(slide, item.get("bullets", []), text_rect, font, int(clamp(body_size, 12, 18)), colors["primary"], colors["accent"])
        if len(images) == 1:
            rects = [image_area]
        else:
            rects = [
                [image_area[0], image_area[1], image_area[2], image_area[3] * 0.48],
                [image_area[0], image_area[1] + image_area[3] * 0.52, image_area[2], image_area[3] * 0.48],
            ]
        for image, rect in zip(images, rects):
            add_rect(slide, rect, colors["surface"], line=True)
            success = add_picture_contained(slide, Path(image["path"]), [rect[0] + 0.06, rect[1] + 0.06, rect[2] - 0.12, rect[3] - 0.12])
            if not success:
                add_textbox(slide, image["name"], [rect[0] + 0.12, rect[1] + 0.2, rect[2] - 0.24, 0.4], font, 10, colors["muted"])
    else:
        render_text_slide(slide, item, profile, colors, preserve, optimize)
        return

    footer_bits = []
    if item.get("documents"):
        footer_bits.append("文档：" + "、".join(item["documents"][:3]))
    if item.get("media"):
        footer_bits.append("媒体：" + "、".join(media["name"] for media in normalize_media_items(item["media"][:2])))
    if footer_bits:
        add_textbox(slide, "；".join(footer_bits), [0.88, slide_h - 0.56, slide_w - 1.76, 0.28], font, 9, colors["muted"])


def render_slide(slide: Any, item: dict[str, Any], profile: dict[str, Any], preserve: set[str], optimize: bool) -> None:
    colors = apply_slide_design(slide, profile, preserve, optimize)
    if item["kind"] == "cover":
        render_cover(slide, item, profile, colors, preserve)
    elif item["kind"] == "section":
        render_section_slide(slide, item, profile, colors, preserve, optimize)
    else:
        render_text_slide(slide, item, profile, colors, preserve, optimize)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def duplicate_slide(prs: Presentation, source_index: int) -> Any:
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)

    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")

    for r_id, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if r_id in dest.part.rels:
            continue
        try:
            dest.part.rels.add_relationship(rel.reltype, rel._target, r_id)
        except AttributeError:
            dest.part.rels._rels[r_id] = _Relationship(
                dest.part.rels._base_uri,
                r_id,
                rel.reltype,
                rel._target_mode,
                rel._target,
            )
        except Exception:
            try:
                dest.part.rels._rels[r_id] = _Relationship(
                    dest.part.rels._base_uri,
                    r_id,
                    rel.reltype,
                    rel._target_mode,
                    rel._target,
                )
            except Exception:
                dest.part.rels._add_relationship(rel.reltype, rel._target, getattr(rel, "is_external", False))
    return dest


def shape_area(shape: Any) -> int:
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def slide_extent(slide: Any) -> tuple[int, int]:
    max_right = 0
    max_bottom = 0
    for shape in slide.shapes:
        try:
            max_right = max(max_right, int(shape.left) + int(shape.width))
            max_bottom = max(max_bottom, int(shape.top) + int(shape.height))
        except Exception:
            continue
    return max(max_right, inches_to_emu(13.333)), max(max_bottom, inches_to_emu(7.5))


def text_shapes(slide: Any) -> list[Any]:
    result: list[Any] = []
    slide_w, slide_h = slide_extent(slide)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        existing_text = getattr(shape, "text", "").strip()
        is_placeholder = bool(getattr(shape, "is_placeholder", False))
        if is_full_slide_shape(shape, slide_w, slide_h) and not existing_text:
            continue
        if existing_text or is_placeholder:
            result.append(shape)
    return result


def first_run_style(shape: Any, fallback_font: str, fallback_size: int, fallback_color: tuple[int, int, int]) -> dict[str, Any]:
    style = {
        "font_name": fallback_font,
        "font_size": fallback_size,
        "bold": False,
        "italic": False,
        "color": fallback_color,
        "alignment": None,
        "level": 0,
        "space_before": None,
        "space_after": None,
        "line_spacing": None,
    }
    try:
        paragraph = shape.text_frame.paragraphs[0]
        style["alignment"] = paragraph.alignment
        style["level"] = paragraph.level
        style["space_before"] = paragraph.space_before
        style["space_after"] = paragraph.space_after
        style["line_spacing"] = paragraph.line_spacing
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                names = run_font_names(run)
                if names:
                    east_asian = next((name for kind, name in names if kind == "east_asian"), None)
                    style["font_name"] = east_asian or names[0][1]
                if run.font.size:
                    style["font_size"] = int(run.font.size.pt)
                if run.font.bold is not None:
                    style["bold"] = bool(run.font.bold)
                if run.font.italic is not None:
                    style["italic"] = bool(run.font.italic)
                rgb = run_rgb(run)
                if rgb:
                    style["color"] = rgb
                if run.text.strip():
                    return style
    except Exception:
        return style
    return style


def apply_run_style(run: Any, style: dict[str, Any]) -> None:
    set_run_font(run, style["font_name"])
    run.font.size = Pt(style["font_size"])
    run.font.bold = style["bold"]
    run.font.italic = style["italic"]
    run.font.color.rgb = tuple_to_rgb(style["color"])


def apply_paragraph_style(paragraph: Any, style: dict[str, Any]) -> None:
    if style.get("alignment") is not None:
        paragraph.alignment = style["alignment"]
    try:
        paragraph.level = int(style.get("level") or 0)
    except Exception:
        pass
    for key in ("space_before", "space_after", "line_spacing"):
        value = style.get(key)
        if value is not None:
            try:
                setattr(paragraph, key, value)
            except Exception:
                pass


def fitted_font_size(shape: Any, lines: list[str], base_size: int) -> int:
    width_in = max(0.6, emu_to_inches(getattr(shape, "width", inches_to_emu(6))))
    height_in = max(0.35, emu_to_inches(getattr(shape, "height", inches_to_emu(1))))
    size = int(base_size)
    text_len = sum(max(1, len(line)) for line in lines)
    hard_lines = max(1, len(lines))
    while size > 9:
        chars_per_line = max(7, int(width_in * 72 / max(7, size) * 1.55))
        visual_lines = sum(max(1, (len(line) + chars_per_line - 1) // chars_per_line) for line in lines)
        max_lines = max(1, int(height_in * 72 / max(8, size * 1.22)))
        if visual_lines <= max_lines + 1 and text_len <= chars_per_line * max_lines * 1.25:
            break
        size -= 1 if hard_lines < 8 else 2
    return max(9, size)


def replace_text_preserve_style(shape: Any, text: str, fallback_font: str, fallback_size: int, fallback_color: tuple[int, int, int]) -> None:
    style = first_run_style(shape, fallback_font, fallback_size, fallback_color)
    lines = [line for line in str(text).splitlines()]
    if not lines:
        lines = [""]
    style["font_size"] = fitted_font_size(shape, lines, int(style["font_size"]))

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        apply_paragraph_style(paragraph, style)
        for run in paragraph.runs:
            apply_run_style(run, style)


def body_text_from_item(item: dict[str, Any]) -> str:
    bullets = item.get("bullets") or []
    if item["kind"] == "agenda":
        return "\n".join(f"{index + 1}. {bullet}" for index, bullet in enumerate(bullets))
    if item["kind"] == "media":
        media = normalize_media_items(item.get("media", []))
        if media:
            label = "音频展示类：" if media_kind_from_path(Path(media[0].get("path") or media[0].get("name") or "")) == "audio" else "产出展示类："
            return label + "\n" + "\n".join(f"{index + 1}. {entry['name']}" for index, entry in enumerate(media[:1]))
    return "\n".join(f"• {bullet}" for bullet in bullets)


def choose_title_and_body_shapes(slide: Any) -> tuple[Any | None, Any | None]:
    shapes = text_shapes(slide)
    if not shapes:
        return None, None
    slide_w, slide_h = slide_extent(slide)
    title_candidates = [
        shape
        for shape in shapes
        if int(shape.top) < slide_h * 0.28 and (max_text_size(shape) >= 22 or int(shape.height) < slide_h * 0.2)
    ]
    title = sorted(title_candidates or shapes, key=lambda sh: (int(sh.top), int(sh.left), -max_text_size(sh)))[0]
    bodies = [shape for shape in shapes if shape is not title]
    body_candidates = [shape for shape in bodies if int(shape.top) > int(title.top) + int(title.height) * 0.25]
    body = max(body_candidates or bodies, key=shape_area) if bodies else None
    return title, body


def add_body_shape_like_template(slide: Any, profile: dict[str, Any]) -> Any:
    rect = body_box(profile)
    return slide.shapes.add_textbox(*rect_to_emu(rect))


def add_title_shape_like_template(slide: Any, profile: dict[str, Any]) -> Any:
    rect = title_box(profile)
    return slide.shapes.add_textbox(*rect_to_emu(rect))


def add_template_images(slide: Any, item: dict[str, Any], profile: dict[str, Any], body_shape: Any | None) -> None:
    images = item.get("images", [])[:2]
    if not images:
        return
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])
    picture_rect = profile["layout"].get("picture_box") or default_profile()["layout"]["picture_box"]
    if body_shape is not None:
        try:
            body_rect = shape_rect_inches(body_shape)
            if body_rect[2] > slide_w * 0.62:
                body_shape.width = Emu(inches_to_emu(body_rect[2] * 0.48))
        except Exception:
            pass
    if picture_rect and picture_rect[2] < slide_w * 0.86 and picture_rect[3] < slide_h * 0.86:
        image_area = [float(picture_rect[0]), float(picture_rect[1]), float(picture_rect[2]), float(picture_rect[3])]
    else:
        image_area = [slide_w * 0.56, 1.35, slide_w * 0.38, slide_h - 2.05]
    rects = [image_area] if len(images) == 1 else [
        [image_area[0], image_area[1], image_area[2], image_area[3] * 0.48],
        [image_area[0], image_area[1] + image_area[3] * 0.52, image_area[2], image_area[3] * 0.48],
    ]
    for image, rect in zip(images, rects):
        try:
            add_picture_contained(slide, Path(image["path"]), rect)
        except Exception:
            continue


def fill_template_slide(slide: Any, item: dict[str, Any], profile: dict[str, Any], log: LogFn, slide_number: int) -> list[dict[str, Any]]:
    font = profile["font"].get("primary") or DEFAULT_FONT
    primary = hex_to_tuple(profile["palette"].get("primary"), (31, 41, 55))
    title_shape, body_shape = choose_title_and_body_shapes(slide)

    if title_shape is None:
        title_shape = add_title_shape_like_template(slide, profile)
    replace_text_preserve_style(title_shape, item.get("title", ""), font, int(profile["font"]["title_size"]), primary)

    if item["kind"] == "cover":
        if body_shape is None:
            body_shape = add_body_shape_like_template(slide, profile)
        replace_text_preserve_style(body_shape, item.get("subtitle", ""), font, int(profile["font"]["subtitle_size"]), primary)
        return []

    text = body_text_from_item(item)
    if body_shape is None:
        body_shape = add_body_shape_like_template(slide, profile)
    replace_text_preserve_style(body_shape, text, font, int(profile["font"]["body_size"]), primary)
    add_template_images(slide, item, profile, body_shape)
    return plan_media_on_slide(slide, item, profile, body_shape, slide_number)


def template_content_slide_indices(prs: Presentation) -> list[int]:
    indices: list[int] = []
    for index, slide in enumerate(prs.slides):
        if index == 0:
            continue
        if len(text_shapes(slide)) >= 2:
            indices.append(index)
    return indices or [min(1, len(prs.slides) - 1)]


def generate_pptx_from_template_pages(
    plan: list[dict[str, Any]],
    profile: dict[str, Any],
    template_path: Path,
    output_pptx: Path,
    log: LogFn,
    embed_media_objects: bool = False,
) -> Path:
    prs = Presentation(str(template_path))
    source_indices = template_content_slide_indices(prs)
    original_count = len(prs.slides)

    while len(prs.slides) < len(plan):
        offset = len(prs.slides) - original_count
        duplicate_slide(prs, source_indices[offset % len(source_indices)])

    while len(prs.slides) > len(plan):
        delete_slide(prs, len(prs.slides) - 1)

    media_insertions: list[dict[str, Any]] = []
    for index, item in enumerate(plan):
        media_insertions.extend(fill_template_slide(prs.slides[index], item, profile, log, index + 1))
        log(f"高保真模板页 {index + 1}/{len(plan)} 已填充：{item['title']}")

    ensure_dir(output_pptx.parent)
    prs.save(str(output_pptx))
    if embed_media_objects:
        embed_media_with_powerpoint(output_pptx, media_insertions, log)
    elif media_insertions:
        log(f"已生成 {len(media_insertions)} 个可点击音视频媒体卡。")
    log(f"PPTX 已保存：{output_pptx}")
    return output_pptx


def powerpoint_pids() -> set[int]:
    try:
        result = subprocess.run(
            ["tasklist", "/FI", "IMAGENAME eq POWERPNT.EXE", "/FO", "CSV", "/NH"],
            capture_output=True,
            text=True,
            timeout=8,
        )
    except Exception:
        return set()
    pids: set[int] = set()
    for line in result.stdout.splitlines():
        parts = [part.strip().strip('"') for part in line.split('","')]
        if len(parts) >= 2 and parts[0].strip('"').lower() == "powerpnt.exe":
            try:
                pids.add(int(parts[1].strip('"')))
            except ValueError:
                continue
    return pids


def kill_processes(pids: set[int]) -> None:
    for pid in pids:
        subprocess.run(["taskkill", "/PID", str(pid), "/F", "/T"], capture_output=True, text=True)


def media_insert_helper(pptx_path: Path, insertions_path: Path) -> int:
    import win32com.client

    insertions = json.loads(insertions_path.read_text(encoding="utf-8"))
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
    try:
        for item in insertions:
            path = Path(item["path"])
            if not path.exists():
                continue
            rect = item["rect"]
            left = float(rect[0]) * POINTS_PER_INCH
            top = float(rect[1]) * POINTS_PER_INCH
            width = float(rect[2]) * POINTS_PER_INCH
            height = float(rect[3]) * POINTS_PER_INCH
            if item.get("kind") == "audio":
                width = min(width, 56)
                height = min(height, 56)
            embed = path.stat().st_size <= MAX_EMBED_MEDIA_BYTES
            shape = presentation.Slides(int(item["slide"])).Shapes.AddMediaObject2(
                str(path),
                not embed,
                embed,
                left,
                top,
                width,
                height,
            )
            shape.Name = f"CodexMedia_{item.get('kind', 'media')}_{path.stem[:24]}"
        presentation.Save()
    finally:
        try:
            presentation.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass
    return 0


def embed_media_with_powerpoint(pptx_path: Path, insertions: list[dict[str, Any]], log: LogFn) -> bool:
    valid_insertions = [item for item in insertions if item.get("path") and Path(item["path"]).exists()]
    insertions = [item for item in valid_insertions if Path(item["path"]).stat().st_size <= MAX_POWERPOINT_INSERT_MEDIA_BYTES]
    skipped = len(valid_insertions) - len(insertions)
    if skipped:
        log(f"{skipped} 个较大音视频已保留为可点击媒体卡，避免 PowerPoint 插入卡死。")
    if not insertions:
        return False
    try:
        with tempfile.TemporaryDirectory(prefix="ppt_media_") as temp_dir:
            temp_dir_path = Path(temp_dir)
            temp_pptx = temp_dir_path / "input.pptx"
            temp_json = temp_dir_path / "media_insertions.json"
            shutil.copy2(pptx_path, temp_pptx)
            temp_json.write_text(json.dumps(insertions, ensure_ascii=False, indent=2), encoding="utf-8")

            before = powerpoint_pids()
            command = [sys.executable, str(Path(__file__).resolve()), "--embed-media-helper", str(temp_pptx), str(temp_json)]
            try:
                result = subprocess.run(command, capture_output=True, text=True, timeout=180)
            except subprocess.TimeoutExpired:
                new_pids = powerpoint_pids() - before
                kill_processes(new_pids)
                log("PowerPoint 插入音视频超过 180 秒，已保留媒体占位卡。")
                return False

            if result.returncode != 0:
                err = (result.stderr or result.stdout or "").strip()
                log(f"PowerPoint 插入音视频失败，已保留媒体占位卡：{err[:2000]}")
                return False
            shutil.copy2(temp_pptx, pptx_path)
        log(f"已插入 {len(insertions)} 个音视频对象。")
        return True
    except Exception as exc:
        log(f"PowerPoint 音视频插入不可用，已保留媒体占位卡：{exc}")
        return False


def export_pdf_with_powerpoint_helper(pptx_path: Path, pdf_path: Path) -> int:
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(pptx_path), WithWindow=False)
    try:
        presentation.SaveAs(str(pdf_path), 32)
    finally:
        try:
            presentation.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass
    return 0


def generate_pptx(
    plan: list[dict[str, Any]],
    profile: dict[str, Any],
    template_path: Path | None,
    output_pptx: Path,
    preserve: set[str],
    optimize: bool,
    log: LogFn,
    embed_media_objects: bool = False,
) -> Path:
    prs = create_presentation(template_path, profile, preserve)
    layout = find_blank_layout(prs)
    media_insertions: list[dict[str, Any]] = []
    for index, item in enumerate(plan, start=1):
        slide = prs.slides.add_slide(layout)
        render_slide(slide, item, profile, preserve, optimize)
        media_insertions.extend(plan_media_on_slide(slide, item, profile, None, index))
        log(f"PPTX 页面 {index}/{len(plan)} 已生成：{item['title']}")
    ensure_dir(output_pptx.parent)
    prs.save(str(output_pptx))
    if embed_media_objects:
        embed_media_with_powerpoint(output_pptx, media_insertions, log)
    elif media_insertions:
        log(f"已生成 {len(media_insertions)} 个可点击音视频媒体卡。")
    log(f"PPTX 已保存：{output_pptx}")
    return output_pptx


def try_export_pdf_with_powerpoint(pptx_path: Path, pdf_path: Path, log: LogFn) -> bool:
    try:
        with tempfile.TemporaryDirectory(prefix="ppt_export_") as temp_dir:
            temp_dir_path = Path(temp_dir)
            temp_pptx = temp_dir_path / "input.pptx"
            temp_pdf = temp_dir_path / "output.pdf"
            shutil.copy2(pptx_path, temp_pptx)
            before_pids = powerpoint_pids()
            command = [sys.executable, str(Path(__file__).resolve()), "--export-pdf-helper", str(temp_pptx), str(temp_pdf)]
            try:
                result = subprocess.run(command, capture_output=True, text=True, timeout=60)
            except subprocess.TimeoutExpired:
                kill_processes(powerpoint_pids() - before_pids)
                log("PowerPoint 导出超过 60 秒，已结束本次导出进程并改用内置 PDF 绘制。")
                return False

            if result.returncode != 0 or not temp_pdf.exists():
                log((result.stdout or result.stderr or "PowerPoint 导出失败").strip())
                kill_processes(powerpoint_pids() - before_pids)
                return False

            ensure_dir(pdf_path.parent)
            shutil.copy2(temp_pdf, pdf_path)
        log(f"PDF 已由 PowerPoint 导出：{pdf_path}")
        return pdf_path.exists()
    except Exception as exc:
        log(f"PowerPoint 导出不可用，改用内置 PDF 绘制：{exc}")
        return False


def pdf_color(color: tuple[int, int, int]) -> tuple[float, float, float]:
    return color[0] / 255, color[1] / 255, color[2] / 255


def rect_points(rect: list[float]) -> fitz.Rect:
    return fitz.Rect(
        rect[0] * POINTS_PER_INCH,
        rect[1] * POINTS_PER_INCH,
        (rect[0] + rect[2]) * POINTS_PER_INCH,
        (rect[1] + rect[3]) * POINTS_PER_INCH,
    )


def wrap_cjk_line(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    lines: list[str] = []
    current = ""
    for char in text:
        current += char
        if len(current) >= max_chars:
            lines.append(current)
            current = ""
    if current:
        lines.append(current)
    return lines


def pdf_add_text(
    page: fitz.Page,
    text: str,
    rect: list[float],
    size: int,
    color: tuple[int, int, int],
    fontfile: str | None,
    align: int = 0,
) -> None:
    fontname = "msyh" if fontfile else "china-s"
    try:
        page.insert_textbox(
            rect_points(rect),
            text,
            fontsize=size,
            fontname=fontname,
            fontfile=fontfile,
            color=pdf_color(color),
            align=align,
        )
    except Exception:
        page.insert_textbox(
            rect_points(rect),
            text,
            fontsize=size,
            fontname="helv",
            color=pdf_color(color),
            align=align,
        )


def pdf_add_rect(page: fitz.Page, rect: list[float], color: tuple[int, int, int], stroke: tuple[int, int, int] | None = None) -> None:
    page.draw_rect(rect_points(rect), color=pdf_color(stroke or color), fill=pdf_color(color), width=0.5 if stroke else 0)


def pdf_add_image(page: fitz.Page, path: Path, rect: list[float]) -> None:
    try:
        page.insert_image(rect_points(rect), filename=str(path), keep_proportion=True)
    except Exception:
        return


def pdf_design(page: fitz.Page, profile: dict[str, Any], preserve: set[str], optimize: bool, colors: dict[str, tuple[int, int, int]]) -> None:
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])
    if "background" in preserve and profile.get("assets", {}).get("background_image"):
        bg = Path(profile["assets"]["background_image"])
        if bg.exists():
            pdf_add_image(page, bg, [0, 0, slide_w, slide_h])
        else:
            pdf_add_rect(page, [0, 0, slide_w, slide_h], colors["background"])
    else:
        pdf_add_rect(page, [0, 0, slide_w, slide_h], colors["background"])
    if optimize:
        if "background" not in preserve:
            pdf_add_rect(page, [0, 0, slide_w, 0.22], colors["accent"])
            pdf_add_rect(page, [0, slide_h - 0.18, slide_w, 0.18], colors["accent_2"])
        pdf_add_rect(page, [0.52, 0.52, 0.08, 5.95], colors["accent"])
    else:
        pdf_add_rect(page, [0.65, 1.22, 1.35, 0.05], colors["accent"])
    if "logo" in preserve:
        for logo in profile.get("assets", {}).get("logos", [])[:3]:
            logo_path = Path(logo["path"])
            if logo_path.exists():
                pdf_add_image(page, logo_path, [float(logo["x"]), float(logo["y"]), float(logo["w"]), float(logo["h"])])


def profile_colors(profile: dict[str, Any], preserve: set[str]) -> dict[str, tuple[int, int, int]]:
    palette = profile["palette"] if "color" in preserve else default_profile()["palette"]
    return {
        "background": hex_to_tuple(palette.get("background"), (247, 248, 251)),
        "surface": hex_to_tuple(palette.get("surface"), (255, 255, 255)),
        "primary": hex_to_tuple(palette.get("primary"), (31, 41, 55)),
        "accent": hex_to_tuple(palette.get("accent"), (47, 128, 237)),
        "accent_2": hex_to_tuple(palette.get("accent_2"), (22, 160, 133)),
        "muted": hex_to_tuple(palette.get("muted"), (107, 114, 128)),
    }


def pdf_render_cover(page: fitz.Page, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], fontfile: str | None, preserve: set[str]) -> None:
    font = profile["font"] if "font" in preserve else default_profile()["font"]
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])
    pdf_add_rect(page, [0.72, 1.28, 0.12, 2.28], colors["accent"])
    pdf_add_text(page, item["title"], [1.05, 1.18, slide_w - 2.0, 1.45], int(font["title_size"]), colors["primary"], fontfile)
    pdf_add_text(page, item["subtitle"], [1.08, 2.75, slide_w - 2.0, 0.55], 15, colors["muted"], fontfile)
    pdf_add_text(page, item["date"], [1.08, 3.36, 4.5, 0.45], 12, colors["muted"], fontfile)
    pdf_add_rect(page, [1.08, slide_h - 1.45, slide_w - 2.16, 0.02], colors["muted"])
    pdf_add_text(page, "由模板学习 + 资料整理自动生成", [1.08, slide_h - 1.18, 5.8, 0.35], 10, colors["muted"], fontfile)


def pdf_render_text(page: fitz.Page, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], fontfile: str | None, preserve: set[str], optimize: bool) -> None:
    font = profile["font"] if "font" in preserve else default_profile()["font"]
    tbox = title_box(profile)
    bbox = body_box(profile)
    pdf_add_text(page, item["title"], tbox, int(clamp(font["title_size"] - 5, 22, 34)), colors["primary"], fontfile)
    if optimize:
        pdf_add_rect(page, [bbox[0] - 0.1, bbox[1] - 0.12, bbox[2] + 0.2, bbox[3] + 0.18], colors["surface"], stroke=darken(colors["surface"], 0.15))
        bbox = [bbox[0] + 0.18, bbox[1] + 0.18, bbox[2] - 0.36, bbox[3] - 0.28]
    max_chars = max(18, int(bbox[2] * 6.6))
    lines: list[str] = []
    for bullet in item.get("bullets", []) or ["暂无正文内容。"]:
        wrapped = wrap_cjk_line(f"• {bullet}", max_chars)
        lines.extend(wrapped)
        lines.append("")
    pdf_add_text(page, "\n".join(lines).strip(), bbox, int(clamp(font["body_size"], 12, 18)), colors["primary"], fontfile)


def pdf_render_section(page: fitz.Page, item: dict[str, Any], profile: dict[str, Any], colors: dict[str, tuple[int, int, int]], fontfile: str | None, preserve: set[str], optimize: bool) -> None:
    images = item.get("images", [])[:2]
    if not images:
        pdf_render_text(page, item, profile, colors, fontfile, preserve, optimize)
        return
    font = profile["font"] if "font" in preserve else default_profile()["font"]
    slide_w = float(profile["slide_width_in"])
    slide_h = float(profile["slide_height_in"])
    pdf_add_text(page, item["title"], title_box(profile), int(clamp(font["title_size"] - 6, 21, 34)), colors["primary"], fontfile)
    text_rect = [0.88, 1.55, slide_w * 0.46, slide_h - 2.1]
    image_area = [slide_w * 0.56, 1.42, slide_w * 0.38, slide_h - 2.05]
    max_chars = max(16, int(text_rect[2] * 5.8))
    lines: list[str] = []
    for bullet in item.get("bullets", []):
        lines.extend(wrap_cjk_line(f"• {bullet}", max_chars))
        lines.append("")
    pdf_add_text(page, "\n".join(lines).strip(), text_rect, int(clamp(font["body_size"], 12, 18)), colors["primary"], fontfile)
    rects = [image_area] if len(images) == 1 else [
        [image_area[0], image_area[1], image_area[2], image_area[3] * 0.48],
        [image_area[0], image_area[1] + image_area[3] * 0.52, image_area[2], image_area[3] * 0.48],
    ]
    for image, rect in zip(images, rects):
        pdf_add_rect(page, rect, colors["surface"], stroke=darken(colors["surface"], 0.15))
        pdf_add_image(page, Path(image["path"]), [rect[0] + 0.06, rect[1] + 0.06, rect[2] - 0.12, rect[3] - 0.12])


def generate_pdf_fallback(plan: list[dict[str, Any]], profile: dict[str, Any], pdf_path: Path, preserve: set[str], optimize: bool, log: LogFn) -> Path:
    ensure_dir(pdf_path.parent)
    doc = fitz.open()
    width = float(profile["slide_width_in"]) * POINTS_PER_INCH
    height = float(profile["slide_height_in"]) * POINTS_PER_INCH
    fontfile = common_windows_cjk_font()
    colors = profile_colors(profile, preserve)
    for index, item in enumerate(plan, start=1):
        page = doc.new_page(width=width, height=height)
        pdf_design(page, profile, preserve, optimize, colors)
        if item["kind"] == "cover":
            pdf_render_cover(page, item, profile, colors, fontfile, preserve)
        elif item["kind"] == "section":
            pdf_render_section(page, item, profile, colors, fontfile, preserve, optimize)
        else:
            pdf_render_text(page, item, profile, colors, fontfile, preserve, optimize)
        log(f"PDF 页面 {index}/{len(plan)} 已绘制：{item['title']}")
    doc.save(pdf_path)
    doc.close()
    log(f"PDF 已保存：{pdf_path}")
    return pdf_path


def generate_outputs(
    template_path: Path,
    source_dir: Path,
    output_dir: Path,
    title: str,
    optimize: bool,
    options: dict[str, Any],
    log: LogFn = noop_log,
) -> dict[str, Path]:
    output_dir = ensure_dir(output_dir.resolve())
    profile = learn_template(template_path, output_dir, log)
    profile = refine_template_with_model(profile, options.get("model", {}), log)
    inventory = scan_source_folder(source_dir, log)
    inventory = refine_with_model(inventory, options.get("model", {}), log)
    preserve = get_preserve_set(options)
    if not preserve:
        preserve = {"page_size"}
    embed_media_objects = bool(options.get("embed_media"))
    plan = build_slide_plan(inventory, profile, title or source_dir.name, optimize, preserve)

    mode_label = "优化设计" if optimize else "同风格"
    base_name = safe_filename(f"{title or source_dir.name}-{mode_label}-{stamp()}")
    plan_path = output_dir / f"{base_name}_slide_plan.json"
    plan_path.write_text(json.dumps({"profile": profile, "inventory": inventory, "slides": plan}, ensure_ascii=False, indent=2), encoding="utf-8")
    log(f"页面计划已保存：{plan_path}")

    pptx_path = output_dir / f"{base_name}.pptx"
    if template_path.suffix.lower() == ".pptx" and not optimize:
        log("使用高保真模板页填充模式：直接复用模板页面骨架。")
        generate_pptx_from_template_pages(plan, profile, template_path, pptx_path, log, embed_media_objects)
    else:
        generate_pptx(plan, profile, template_path if template_path.suffix.lower() == ".pptx" else None, pptx_path, preserve, optimize, log, embed_media_objects)
    return {"pptx": pptx_path, "plan": plan_path}


def markdown_source_preview(inventory: dict[str, Any]) -> str:
    lines = [
        "# 资料摘要",
        "",
        f"- 资料目录：{inventory['source_dir']}",
        f"- 分组数：{inventory['stats']['sections']}",
        f"- 文档数：{inventory['stats']['documents']}",
        f"- 图片数：{inventory['stats']['images']}",
        f"- 音视频数：{inventory['stats']['media']}",
        "",
    ]
    for section in inventory["sections"]:
        lines.append(f"## {section['title']}")
        for bullet in section.get("bullets", [])[:8]:
            lines.append(f"- {bullet}")
        if section.get("documents"):
            docs = "、".join(doc["relative"] for doc in section["documents"][:8])
            lines.append(f"- 文档：{docs}")
        if section.get("images"):
            images = "、".join(image["relative"] for image in section["images"][:8])
            lines.append(f"- 图片：{images}")
        if section.get("media"):
            media = "、".join(item["relative"] for item in section["media"][:8])
            lines.append(f"- 媒体：{media}")
        lines.append("")
    if inventory.get("warnings"):
        lines.extend(["## 解析提醒", ""])
        for warning in inventory["warnings"]:
            lines.append(f"- {warning}")
    return "\n".join(lines).strip() + "\n"


def codex_prompt_text(task: dict[str, Any], task_path: Path) -> str:
    mode = "优化/新设计" if task["optimize"] else "同风格高保真"
    preserve = "、".join(key for key, value in task["preserve"].items() if value) or "无"
    return f"""请处理这个 PPTX 生成任务包：

{task_path}

我的目标：
1. 读取 codex_task.json、模板文件、资料文件夹和资料摘要。
2. 模板路径：{task['template_path']}
3. 资料路径：{task['source_dir']}
4. 输出目录：{task['output_dir']}
5. 报告标题：{task['title']}
6. 生成模式：{mode}
7. 需要保留的模板模块：{preserve}
8. 真实音视频插入：{"启用" if task.get("embed_media") else "不启用，默认使用可点击媒体卡"}

请你直接分析模板 PPTX 的真实视觉风格、页面结构、字体、配色、留白、图片处理方式、音视频展示位置和内容节奏，然后生成新的 PPTX。
如果目标是同风格，请优先复用模板页面骨架，避免重新画一套不像的版式。
如果目标是优化/新设计，请先说明你会保留哪些模板特征，再生成结果。
完成后请告诉我生成文件路径，并检查 PPTX 页数、字体、版式、图片和音视频位置是否正常。
"""


def create_codex_task_package(
    template_path: Path,
    source_dir: Path,
    output_dir: Path,
    title: str,
    optimize: bool,
    options: dict[str, Any],
    log: LogFn = noop_log,
) -> dict[str, Path]:
    output_dir = ensure_dir(output_dir.resolve())
    task_root = ensure_dir(output_dir / "Codex任务包")
    task_dir = ensure_dir(task_root / safe_filename(f"{title or source_dir.name}-{stamp()}", "Codex任务包"))

    log("开始生成 Codex 任务包。")
    profile = learn_template(template_path, task_dir, log)
    inventory = scan_source_folder(source_dir, log)
    preserve = get_preserve_set(options)
    if not preserve:
        preserve = {"page_size"}
    plan = build_slide_plan(inventory, profile, title or source_dir.name, optimize, preserve)

    profile_path = task_dir / "template_profile.json"
    inventory_path = task_dir / "source_inventory.json"
    plan_path = task_dir / "slide_plan.json"
    preview_path = task_dir / "资料摘要.md"
    task_path = task_dir / "codex_task.json"
    prompt_path = task_dir / "给Codex的提示.txt"

    profile_path.write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
    inventory_path.write_text(json.dumps(inventory, ensure_ascii=False, indent=2), encoding="utf-8")
    plan_path.write_text(json.dumps({"slides": plan}, ensure_ascii=False, indent=2), encoding="utf-8")
    preview_path.write_text(markdown_source_preview(inventory), encoding="utf-8")

    task = {
        "version": 1,
        "created_at": dt.datetime.now().isoformat(timespec="seconds"),
        "workspace": str(Path.cwd()),
        "template_path": str(template_path.resolve()),
        "source_dir": str(source_dir.resolve()),
        "output_dir": str(output_dir.resolve()),
        "title": title or source_dir.name,
        "optimize": optimize,
        "preserve": {key: key in preserve for key in default_profile()["modules"].keys()},
        "embed_media": bool(options.get("embed_media")),
        "task_dir": str(task_dir),
        "profile_path": str(profile_path),
        "inventory_path": str(inventory_path),
        "slide_plan_path": str(plan_path),
        "source_preview_path": str(preview_path),
        "expected_outputs": ["pptx"],
        "notes": [
            "这个任务包用于交给当前 Codex 会话处理，不会让本地工具直接调用 ChatGPT/Codex 订阅。",
            "Codex 可读取这些路径并继续生成或优化 PPTX。",
        ],
    }
    task_path.write_text(json.dumps(task, ensure_ascii=False, indent=2), encoding="utf-8")
    prompt = codex_prompt_text(task, task_path)
    prompt_path.write_text(prompt, encoding="utf-8")

    latest_path = WORK_DIR / "最新Codex任务包.txt"
    latest_path.write_text(str(task_path), encoding="utf-8")
    log(f"Codex 任务包已生成：{task_path}")
    log(f"给 Codex 的提示已生成：{prompt_path}")
    return {
        "task": task_path,
        "prompt": prompt_path,
        "profile": profile_path,
        "inventory": inventory_path,
        "plan": plan_path,
        "preview": preview_path,
        "latest": latest_path,
    }


def generate_pair_mapped_ppt(
    learn_ppt: Path,
    source_dir: Path,
    output_dir: Path,
    title: str,
    learn_dir: Path | None = None,
    log: LogFn = noop_log,
) -> dict[str, Path]:
    if not PAIR_PACKAGE_SCRIPT.exists():
        raise FileNotFoundError(f"样本映射脚本不存在：{PAIR_PACKAGE_SCRIPT}")
    if not learn_ppt.exists() or learn_ppt.suffix.lower() != ".pptx":
        raise FileNotFoundError(f"学习PPT不存在或不是 .pptx：{learn_ppt}")
    if not source_dir.exists() or not source_dir.is_dir():
        raise FileNotFoundError(f"新资料夹不存在：{source_dir}")

    learn_dir = learn_dir or learn_ppt.parent
    output_dir = ensure_dir(output_dir)
    command = [
        sys.executable,
        str(PAIR_PACKAGE_SCRIPT),
        "--root",
        str(PROJECT_ROOT),
        "--learn-dir",
        str(learn_dir),
        "--learn-ppt",
        str(learn_ppt),
        "--source-dir",
        str(source_dir),
        "--output",
        str(output_dir),
        "--title",
        title or f"{source_dir.name} - 样本映射媒体直替版",
    ]
    log("启动样本映射直替流程：学习旧资料夹和旧成品PPT，再替换为新资料夹内容。")
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    result = subprocess.run(
        command,
        cwd=str(PROJECT_ROOT),
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="backslashreplace",
        env=env,
        timeout=60 * 60,
    )
    for line in (result.stdout or "").splitlines():
        log(line)
    for line in (result.stderr or "").splitlines():
        log(line)
    if result.returncode != 0:
        raise RuntimeError(f"样本映射生成失败，退出码 {result.returncode}")

    outputs: dict[str, Path] = {}
    for line in (result.stdout or "").splitlines():
        if line.startswith("PPTX="):
            outputs["pptx"] = Path(line.split("=", 1)[1].strip())
        elif line.startswith("PLAN="):
            outputs["plan"] = Path(line.split("=", 1)[1].strip())
    if "pptx" not in outputs:
        raise RuntimeError("样本映射脚本没有返回 PPTX 输出路径。")
    if "plan" not in outputs:
        plan_guess = outputs["pptx"].with_name(outputs["pptx"].stem + "_package_plan.json")
        if plan_guess.exists():
            outputs["plan"] = plan_guess
    return outputs


def profile_summary(profile: dict[str, Any]) -> str:
    font = profile["font"]
    palette = profile["palette"]
    modules = profile["modules"]
    analysis = profile["analysis"]
    lines = [
        f"模板类型：{profile['template_type']}",
        f"页面尺寸：{profile['slide_width_in']} x {profile['slide_height_in']} 英寸",
        f"页数/版式：{analysis.get('slide_count', 0)} / {analysis.get('layout_count', 0)}",
        f"字体：{font['primary']}，标题 {font['title_size']}pt，正文 {font['body_size']}pt",
        "颜色：" + "、".join(f"{k}={v}" for k, v in palette.items()),
        "可保留模块：" + "、".join(k for k, v in modules.items() if v),
    ]
    notes = analysis.get("notes") or []
    lines.extend(notes)
    return "\n".join(lines)


class GeneratorApp:
    def __init__(self, root: tk.Tk) -> None:
        self.root = root
        self.root.title(APP_NAME)
        self.root.geometry("980x720")
        self.template_var = tk.StringVar(value="")
        self.source_var = tk.StringVar(value="")
        self.output_var = tk.StringVar(value=str(DEFAULT_OUTPUT_DIR))
        self.title_var = tk.StringVar(value="")
        self.model_mode_var = tk.StringVar(value="规则引擎")
        self.model_name_var = tk.StringVar(value="")
        self.base_url_var = tk.StringVar(value="")
        self.api_key_var = tk.StringVar(value="")
        self.embed_media_var = tk.BooleanVar(value=False)
        self.profile: dict[str, Any] | None = None
        self.preserve_vars: dict[str, tk.BooleanVar] = {
            "font": tk.BooleanVar(value=True),
            "color": tk.BooleanVar(value=True),
            "background": tk.BooleanVar(value=True),
            "logo": tk.BooleanVar(value=True),
            "layout": tk.BooleanVar(value=True),
            "page_size": tk.BooleanVar(value=True),
            "picture_style": tk.BooleanVar(value=True),
        }
        self._build()

    def _build(self) -> None:
        root = self.root
        root.columnconfigure(0, weight=1)
        root.rowconfigure(4, weight=1)

        top = ttk.Frame(root, padding=14)
        top.grid(row=0, column=0, sticky="ew")
        top.columnconfigure(1, weight=1)

        self._row_picker(top, 0, "学习PPT/模板", self.template_var, self.pick_template, "推荐选择上一期成品 PPTX")
        self._row_picker(top, 1, "资料文件夹", self.source_var, self.pick_source, "选择要生成报告的资料文件夹")
        self._row_picker(top, 2, "输出目录", self.output_var, self.pick_output, "选择生成结果保存位置")

        ttk.Label(top, text="报告标题").grid(row=3, column=0, sticky="w", pady=5)
        ttk.Entry(top, textvariable=self.title_var).grid(row=3, column=1, sticky="ew", padx=8, pady=5)

        preserve = ttk.LabelFrame(root, text="优化/新设计时保留的模板模块", padding=12)
        preserve.grid(row=1, column=0, sticky="ew", padx=14, pady=(0, 8))
        labels = {
            "font": "字体",
            "color": "配色",
            "background": "底图/背景",
            "logo": "Logo/水印",
            "layout": "母版/版式",
            "page_size": "页面比例",
            "picture_style": "配图比例",
        }
        for index, (key, label) in enumerate(labels.items()):
            ttk.Checkbutton(preserve, text=label, variable=self.preserve_vars[key]).grid(row=0, column=index, sticky="w", padx=(0, 14))

        model = ttk.LabelFrame(root, text="模型设置（可选，不填则使用规则引擎）", padding=12)
        model.grid(row=2, column=0, sticky="ew", padx=14, pady=(0, 8))
        model.columnconfigure(3, weight=1)
        ttk.Label(model, text="模式").grid(row=0, column=0, sticky="w")
        combo = ttk.Combobox(
            model,
            textvariable=self.model_mode_var,
            values=["规则引擎", "Ollama 本地模型", "OpenAI 兼容远程模型"],
            width=20,
            state="readonly",
        )
        combo.grid(row=0, column=1, sticky="w", padx=(6, 18))
        ttk.Label(model, text="模型名").grid(row=0, column=2, sticky="w")
        ttk.Entry(model, textvariable=self.model_name_var, width=22).grid(row=0, column=3, sticky="w", padx=(6, 18))
        ttk.Label(model, text="Base URL").grid(row=1, column=0, sticky="w", pady=(8, 0))
        ttk.Entry(model, textvariable=self.base_url_var).grid(row=1, column=1, columnspan=3, sticky="ew", padx=(6, 18), pady=(8, 0))
        ttk.Label(model, text="API Key").grid(row=1, column=4, sticky="w", pady=(8, 0))
        ttk.Entry(model, textvariable=self.api_key_var, show="*", width=28).grid(row=1, column=5, sticky="w", padx=(6, 0), pady=(8, 0))
        ttk.Checkbutton(model, text="插入真实音视频对象", variable=self.embed_media_var).grid(row=2, column=0, columnspan=6, sticky="w", pady=(8, 0))

        actions = ttk.Frame(root, padding=(14, 0, 14, 10))
        actions.grid(row=3, column=0, sticky="ew")
        ttk.Button(actions, text="分析模板", command=self.analyze_template).pack(side="left", padx=(0, 8))
        ttk.Button(actions, text="样本映射直替PPT（推荐）", command=self.run_pair_mapped_generation).pack(side="left", padx=(0, 8))
        ttk.Button(actions, text="生成同风格（高保真）PPT", command=lambda: self.run_generation(False)).pack(side="left", padx=(0, 8))
        ttk.Button(actions, text="优化/新设计", command=lambda: self.run_generation(True)).pack(side="left", padx=(0, 8))
        ttk.Button(actions, text="交给 Codex 处理", command=self.create_codex_task).pack(side="left", padx=(0, 8))
        ttk.Button(actions, text="打开输出目录", command=self.open_output_dir).pack(side="right")

        log_frame = ttk.LabelFrame(root, text="运行日志", padding=8)
        log_frame.grid(row=4, column=0, sticky="nsew", padx=14, pady=(0, 14))
        log_frame.rowconfigure(0, weight=1)
        log_frame.columnconfigure(0, weight=1)
        self.log_text = tk.Text(log_frame, height=18, wrap="word")
        self.log_text.grid(row=0, column=0, sticky="nsew")
        scrollbar = ttk.Scrollbar(log_frame, command=self.log_text.yview)
        scrollbar.grid(row=0, column=1, sticky="ns")
        self.log_text.configure(yscrollcommand=scrollbar.set)
        self.log(f"{APP_NAME} 已启动。正式周报建议选择上一期成品 PPT，并点击“样本映射直替PPT（推荐）”。")

    def _row_picker(self, parent: ttk.Frame, row: int, label: str, variable: tk.StringVar, command: Callable[[], None], hint: str) -> None:
        ttk.Label(parent, text=label).grid(row=row, column=0, sticky="w", pady=5)
        ttk.Entry(parent, textvariable=variable).grid(row=row, column=1, sticky="ew", padx=8, pady=5)
        ttk.Button(parent, text="选择", command=command).grid(row=row, column=2, sticky="e", pady=5)
        ttk.Label(parent, text=hint, foreground="#666666").grid(row=row, column=3, sticky="w", padx=(8, 0), pady=5)

    def log(self, message: str) -> None:
        def append() -> None:
            time = dt.datetime.now().strftime("%H:%M:%S")
            self.log_text.insert("end", f"[{time}] {message}\n")
            self.log_text.see("end")

        self.root.after(0, append)

    def show_info(self, message: str) -> None:
        self.root.after(0, lambda: messagebox.showinfo(APP_NAME, message))

    def show_error(self, message: str) -> None:
        self.root.after(0, lambda: messagebox.showerror(APP_NAME, message))

    def pick_template(self) -> None:
        path = filedialog.askopenfilename(
            title="选择学习模板",
            filetypes=[("PowerPoint", "*.pptx")],
        )
        if path:
            self.template_var.set(path)
            if not self.title_var.get().strip():
                self.title_var.set(Path(path).stem)

    def pick_source(self) -> None:
        path = filedialog.askdirectory(title="选择资料文件夹")
        if path:
            self.source_var.set(path)
            if not self.title_var.get().strip():
                self.title_var.set(Path(path).name)

    def pick_output(self) -> None:
        path = filedialog.askdirectory(title="选择输出目录")
        if path:
            self.output_var.set(path)

    def open_output_dir(self) -> None:
        output = ensure_dir(Path(self.output_var.get() or DEFAULT_OUTPUT_DIR))
        os.startfile(str(output))

    def validate_paths(self) -> tuple[Path, Path, Path] | None:
        template = Path(self.template_var.get().strip('" '))
        source = Path(self.source_var.get().strip('" '))
        output = Path(self.output_var.get().strip('" ') or DEFAULT_OUTPUT_DIR)
        if not template.exists() or template.suffix.lower() not in SUPPORTED_TEMPLATE_EXTS:
            messagebox.showerror(APP_NAME, "请先选择有效的 .pptx 学习PPT或模板。")
            return None
        if not source.exists() or not source.is_dir():
            messagebox.showerror(APP_NAME, "请先选择有效的资料文件夹。")
            return None
        return template, source, output

    def analyze_template(self) -> None:
        valid = self.validate_paths()
        if not valid:
            return
        template, _, output = valid

        def worker() -> None:
            try:
                profile = learn_template(template, output, self.log)
                self.profile = profile
                self.log("模板分析摘要：\n" + profile_summary(profile))
            except Exception:
                self.log(traceback.format_exc())
                self.show_error("模板分析失败，详情请查看日志。")

        threading.Thread(target=worker, daemon=True).start()

    def options(self) -> dict[str, Any]:
        return {
            "preserve": {key: var.get() for key, var in self.preserve_vars.items()},
            "model": {
                "mode": self.model_mode_var.get(),
                "model": self.model_name_var.get(),
                "base_url": self.base_url_var.get(),
                "api_key": self.api_key_var.get(),
            },
            "embed_media": self.embed_media_var.get(),
        }

    def run_generation(self, optimize: bool) -> None:
        valid = self.validate_paths()
        if not valid:
            return
        template, source, output = valid
        title = self.title_var.get().strip() or source.name

        def worker() -> None:
            try:
                self.log("开始生成优化设计版。" if optimize else "开始生成同风格版。")
                result = generate_outputs(template, source, output, title, optimize, self.options(), self.log)
                self.log("生成完成：")
                for key, path in result.items():
                    self.log(f"{key.upper()}：{path}")
                self.show_info(f"生成完成：\n{result['pptx']}")
            except Exception:
                self.log(traceback.format_exc())
                self.show_error("生成失败，详情请查看日志。")

        threading.Thread(target=worker, daemon=True).start()

    def run_pair_mapped_generation(self) -> None:
        valid = self.validate_paths()
        if not valid:
            return
        learn_ppt, source, output = valid
        title = self.title_var.get().strip() or f"{source.name} - 样本映射媒体直替版"

        def worker() -> None:
            try:
                result = generate_pair_mapped_ppt(learn_ppt, source, output, title, learn_ppt.parent, self.log)
                self.log("样本映射生成完成：")
                for key, path in result.items():
                    self.log(f"{key.upper()}：{path}")
                self.show_info(f"样本映射生成完成：\n{result['pptx']}")
            except Exception:
                self.log(traceback.format_exc())
                self.show_error("样本映射生成失败，详情请查看日志。")

        threading.Thread(target=worker, daemon=True).start()

    def create_codex_task(self) -> None:
        valid = self.validate_paths()
        if not valid:
            return
        template, source, output = valid
        title = self.title_var.get().strip() or source.name

        def worker() -> None:
            try:
                result = create_codex_task_package(template, source, output, title, False, self.options(), self.log)
                prompt = result["prompt"].read_text(encoding="utf-8")

                def copy_and_show() -> None:
                    try:
                        self.root.clipboard_clear()
                        self.root.clipboard_append(prompt)
                    except Exception:
                        pass
                    messagebox.showinfo(
                        APP_NAME,
                        "Codex 任务包已生成，提示词已复制到剪贴板：\n"
                        f"{result['task']}\n\n"
                        "回到 Codex 对话里粘贴，或直接说：处理最新任务包。",
                    )

                self.root.after(0, copy_and_show)
            except Exception:
                self.log(traceback.format_exc())
                self.show_error("Codex 任务包生成失败，详情请查看日志。")

        threading.Thread(target=worker, daemon=True).start()


def parse_cli(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=APP_NAME)
    parser.add_argument("--template", required=True, help="学习模板路径，支持 .pptx")
    parser.add_argument("--source", required=True, help="资料文件夹")
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT_DIR), help="输出目录")
    parser.add_argument("--title", default="", help="报告标题")
    parser.add_argument("--optimize", action="store_true", help="生成优化设计版")
    parser.add_argument("--codex-task", action="store_true", help="只生成 Codex 任务包，不直接生成 PPT")
    parser.add_argument("--pair-map", action="store_true", help="使用旧资料夹+旧成品PPT样本映射直替流程（推荐）")
    parser.add_argument("--learn-dir", default="", help="样本映射学习资料夹；不填则使用 --template 所在目录")
    parser.add_argument(
        "--preserve",
        default="font,color,background,logo,layout,page_size,picture_style",
        help="优化时保留模块，逗号分隔：font,color,background,logo,layout,page_size,picture_style",
    )
    parser.add_argument("--model-mode", default="规则引擎", choices=["规则引擎", "Ollama 本地模型", "OpenAI 兼容远程模型"])
    parser.add_argument("--model", default="")
    parser.add_argument("--base-url", default="")
    parser.add_argument("--api-key", default="")
    parser.add_argument("--embed-media", action="store_true", help="尝试用 PowerPoint 插入真实音视频对象；默认只生成可点击媒体卡")
    return parser.parse_args(argv)


def cli_main(argv: list[str]) -> int:
    args = parse_cli(argv)
    preserve_keys = {item.strip() for item in args.preserve.split(",") if item.strip()}
    options = {
        "preserve": {
            "font": "font" in preserve_keys,
            "color": "color" in preserve_keys,
            "background": "background" in preserve_keys,
            "logo": "logo" in preserve_keys,
            "layout": "layout" in preserve_keys,
            "page_size": "page_size" in preserve_keys,
            "picture_style": "picture_style" in preserve_keys,
        },
        "model": {
            "mode": args.model_mode,
            "model": args.model,
            "base_url": args.base_url,
            "api_key": args.api_key,
        },
        "embed_media": args.embed_media,
    }

    def log(message: str) -> None:
        safe_console_print(message)

    if args.pair_map:
        result = generate_pair_mapped_ppt(
            Path(args.template),
            Path(args.source),
            Path(args.output),
            args.title or f"{Path(args.source).name} - 样本映射媒体直替版",
            Path(args.learn_dir) if args.learn_dir else Path(args.template).parent,
            log,
        )
        safe_console_print(json.dumps({key: str(value) for key, value in result.items()}, ensure_ascii=False, indent=2))
        return 0

    if args.codex_task:
        result = create_codex_task_package(
            Path(args.template),
            Path(args.source),
            Path(args.output),
            args.title,
            args.optimize,
            options,
            log,
        )
        safe_console_print(json.dumps({key: str(value) for key, value in result.items()}, ensure_ascii=False, indent=2))
        return 0

    result = generate_outputs(
        Path(args.template),
        Path(args.source),
        Path(args.output),
        args.title,
        args.optimize,
        options,
        log,
    )
    safe_console_print(json.dumps({key: str(value) for key, value in result.items()}, ensure_ascii=False, indent=2))
    return 0


def gui_main() -> int:
    if tk is None:
        print("当前 Python 环境无法加载 tkinter，请使用命令行模式。", file=sys.stderr)
        return 2
    root = tk.Tk()
    try:
        style = ttk.Style(root)
        if "vista" in style.theme_names():
            style.theme_use("vista")
    except Exception:
        pass
    GeneratorApp(root)
    root.mainloop()
    return 0


def main() -> int:
    if len(sys.argv) >= 4 and sys.argv[1] == "--embed-media-helper":
        return media_insert_helper(Path(sys.argv[2]), Path(sys.argv[3]))
    if len(sys.argv) >= 4 and sys.argv[1] == "--export-pdf-helper":
        return export_pdf_with_powerpoint_helper(Path(sys.argv[2]), Path(sys.argv[3]))
    if len(sys.argv) > 1:
        return cli_main(sys.argv[1:])
    return gui_main()


if __name__ == "__main__":
    raise SystemExit(main())
