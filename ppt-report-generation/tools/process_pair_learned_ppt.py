from __future__ import annotations

import datetime as dt
import json
import os
import re
import shutil
import tempfile
import time
import xml.etree.ElementTree as ET
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


ROOT = Path.cwd()
DEFAULT_LEARN_DIR_NAME = "2026 0518 - 0529 -  双周报"
DEFAULT_SOURCE_DIR_NAME = "2026 0601 - 0612 -  双周报"
DEFAULT_LEARN_PPT_NAME = "音频中心工作汇报 0518 - 0529.pptx"


def safe_filename(text: str, default: str = "样本学习直插版") -> str:
    cleaned = re.sub(r'[<>:"/\\|?*\r\n\t]+', "_", str(text)).strip(" ._")
    return cleaned[:100] or default


def path_from_value(value: str | os.PathLike | None, base: Path) -> Path | None:
    if value is None or str(value).strip() == "":
        return None
    path = Path(str(value).strip('" '))
    return path if path.is_absolute() else base / path


def find_learn_ppt(learn_dir: Path) -> Path:
    direct_default = learn_dir / DEFAULT_LEARN_PPT_NAME
    if direct_default.exists():
        return direct_default
    candidates = [
        path
        for path in learn_dir.glob("*.pptx")
        if path.is_file() and not path.name.startswith("~$")
    ]
    if not candidates:
        return direct_default
    candidates.sort(
        key=lambda path: (
            "工作汇报" not in path.stem,
            "生成" in path.stem or "预览" in path.stem,
            path.name.lower(),
        )
    )
    return candidates[0]


def configure_paths(
    root: str | os.PathLike | None = None,
    learn_dir: str | os.PathLike | None = None,
    source_dir: str | os.PathLike | None = None,
    learn_ppt: str | os.PathLike | None = None,
    output_dir: str | os.PathLike | None = None,
    title: str | None = None,
) -> None:
    global ROOT, LEARN_DIR, SOURCE_DIR, LEARN_PPT, OUTPUT_DIR, STAMP, BASE_NAME, PPTX_OUT, PLAN_OUT

    cwd = Path.cwd()
    ROOT = path_from_value(root or os.environ.get("PPT_PAIR_ROOT"), cwd) or cwd
    LEARN_DIR = path_from_value(learn_dir or os.environ.get("PPT_PAIR_LEARN_DIR"), ROOT) or ROOT / DEFAULT_LEARN_DIR_NAME
    SOURCE_DIR = path_from_value(source_dir or os.environ.get("PPT_PAIR_SOURCE_DIR"), ROOT) or ROOT / DEFAULT_SOURCE_DIR_NAME
    LEARN_PPT = path_from_value(learn_ppt or os.environ.get("PPT_PAIR_LEARN_PPT"), ROOT) or find_learn_ppt(LEARN_DIR)
    OUTPUT_DIR = path_from_value(output_dir or os.environ.get("PPT_PAIR_OUTPUT_DIR"), ROOT) or ROOT / "PDF生成工具" / "生成结果"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    STAMP = dt.datetime.now().strftime("%Y%m%d-%H%M%S")
    base_title = title or os.environ.get("PPT_PAIR_TITLE") or f"{SOURCE_DIR.name} - 样本学习直插版"
    BASE_NAME = f"{safe_filename(base_title)}-{STAMP}"
    PPTX_OUT = OUTPUT_DIR / f"{BASE_NAME}.pptx"
    PLAN_OUT = OUTPUT_DIR / f"{BASE_NAME}_pair_plan.json"


configure_paths()

W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
MEDIA_EXTS = {".mp4", ".mov", ".m4v", ".avi", ".wmv", ".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}
VIDEO_EXTS = {".mp4", ".mov", ".m4v", ".avi", ".wmv"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}
IGNORE_PARTS = {"AI修正版预览", "AI生成预览", "html_report_0518_0529", "_html_work", "最终修改版预览", "_internal"}
EMU_PER_POINT = 12700


def read_docx_lines(path: Path) -> list[str]:
    try:
        with ZipFile(path) as archive:
            xml = archive.read("word/document.xml")
    except Exception:
        return []
    root = ET.fromstring(xml)
    lines: list[str] = []
    for paragraph in root.iter(W_NS + "p"):
        text = "".join(node.text or "" for node in paragraph.iter(W_NS + "t")).strip()
        if text:
            lines.append(" ".join(text.split()))
    compact: list[str] = []
    for line in lines:
        if not compact or compact[-1] != line:
            compact.append(line)
    return compact


def read_txt(path: Path) -> list[str]:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return []
    return [line.strip() for line in text.splitlines() if line.strip()]


def doc_lines(relative: str) -> list[str]:
    return read_docx_lines(SOURCE_DIR / relative)


def txt_lines(relative: str) -> list[str]:
    return read_txt(SOURCE_DIR / relative)


def media_kind(path: Path) -> str:
    if path.suffix.lower() in AUDIO_EXTS:
        return "audio"
    if path.suffix.lower() in VIDEO_EXTS:
        return "video"
    return "media"


def existing(paths: list[Path]) -> list[Path]:
    return [path for path in paths if path.exists()]


def list_media_under(relative: str) -> list[Path]:
    root = SOURCE_DIR / relative
    if not root.exists():
        return []
    items = []
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in MEDIA_EXTS:
            continue
        rel_parts = set(path.relative_to(SOURCE_DIR).parts)
        if rel_parts & IGNORE_PARTS:
            continue
        items.append(path)
    return sorted(items, key=lambda item: str(item.relative_to(SOURCE_DIR)).lower())


def media_groups() -> dict[str, list[Path]]:
    def ordered_media(relative: str, preferred: list[Path]) -> list[Path]:
        result = existing(preferred)
        seen = {str(path).lower() for path in result}
        for path in list_media_under(relative):
            key = str(path).lower()
            if key not in seen:
                result.append(path)
                seen.add(key)
        return result

    return {
        "L": ordered_media(
            "L",
            [
                SOURCE_DIR / "L" / "UI.mp4",
                SOURCE_DIR / "L" / "loading.mp4",
                SOURCE_DIR / "L" / "郑青" / "试听.mp4",
            ],
        ),
        "LL": ordered_media(
            "LL",
            [
                SOURCE_DIR / "LL" / "安家辉" / "HO 10201 Combat Skill021.mp4",
                SOURCE_DIR / "LL" / "安家辉" / "HO 10241 Cutscene ExSkill.mp4",
                SOURCE_DIR / "LL" / "安家辉" / "HO 10251 Combat CoopSkill Sherry.mp4",
                SOURCE_DIR / "LL" / "安家辉" / "HO 10251 Combat CoopSkill Sherry TL.mp4",
                SOURCE_DIR / "LL" / "朱静雅" / "ElevenLabs视频配乐演示.mp4",
                SOURCE_DIR / "LL" / "朱静雅" / "ElevenLabs语音演示.mp4",
                SOURCE_DIR / "LL" / "朱静雅" / "ElevenLabs音效演示.mp4",
                SOURCE_DIR / "LL" / "朱静雅" / "Sonilo演示.mp4",
                SOURCE_DIR / "LL" / "朱静雅" / "Suno生成.mp3",
                SOURCE_DIR / "LL" / "吉禹鹏" / "飞书20260612-165541.mp4",
            ],
        ),
        "U": ordered_media(
            "U",
            [
                SOURCE_DIR / "U" / "郑青" / "试听.mp4",
                SOURCE_DIR / "U" / "郑青" / "BGM_Theme_DuWanSanQing.wav",
            ],
        ),
        "FM": ordered_media("FM", [SOURCE_DIR / "FM" / "冯博宇.mp4"]),
        "EF": ordered_media(
            "EF",
            [
                SOURCE_DIR / "EF" / "0610 UI替换.mp4",
                SOURCE_DIR / "EF" / "Stamina.mp4",
                SOURCE_DIR / "EF" / "成就.mp4",
            ],
        ),
        "TA": ordered_media(
            "音频技术",
            [SOURCE_DIR / "音频技术" / "关晗怡 联调工具" / "联调工具 显示特效 动画 timeline，点击定位到timeline.mp4"],
        ),
    }


def source_text_plan() -> dict[int, dict[str, str]]:
    l_lines = doc_lines(r"L\L.docx") + doc_lines(r"L\郑青\L.docx")
    ll_lines = doc_lines(r"LL\关晗怡.docx")
    u_lines = doc_lines(r"U\郑青\U.docx")
    fm_lines = txt_lines(r"FM\冯博宇.txt")
    tech_feng = doc_lines(r"音频技术\冯博宇.docx")
    tech_cai = doc_lines(r"音频技术\蔡新星.docx")
    tech_zheng = doc_lines(r"音频技术\郑青.docx")

    def numbered(lines: list[str], limit: int = 5) -> str:
        return "\r".join(f"{index}. {line}" for index, line in enumerate(lines[:limit], 1))

    return {
        1: {"（5月18日-5月29日 周报）": "（6月1日-6月12日 周报）"},
        2: {
            "音频中台 |": "                  音频中台 | 虞鹏 | 郑青 | 蔡新星 | 安家辉 |",
            "| 冯博宇": "冯博宇 | 朱静雅 | 关晗怡 | 吴帅斐 | 高世佳",
            "吉禹鹏": "冯博宇 | 朱静雅 | 关晗怡 | 吴帅斐 | 高世佳",
        },
        3: {
            "参与人员": "参与人员\r音效设计：郑青\r音频技术：冯博宇，蔡新星\r本期重点：0618、0716版本对接，语音标准化"
        },
        4: {"设计&沟通类：": "设计&沟通类：\r1. 0618、0716版本对接\r2. 语音标准化\r3. 版本资源跟进与制作排期同步"},
        5: {"资源制作类：": "资源制作类：\r" + numbered(l_lines, 4)},
        8: {
            "参与人员": "参与人员\r音效设计：关晗怡，安家辉，吴帅斐，虞鹏\r音乐设计：朱静雅，吉禹鹏，虞鹏\r音频策划：关晗怡，虞鹏\r音频技术：冯博宇，蔡新星"
        },
        9: {"设计&沟通类：": "设计&沟通类：\r" + numbered(ll_lines[:6], 6)},
        10: {
            "资源制作 / 问题处理类：": "资源制作 / 问题处理类：\r1. 战斗技能、演出技能与视频展示\r2. AI工具制作与工具需求梳理\r3. 音效缺漏补齐日志更新\r4. Prefab / Spine / Timeline 问题排查"
        },
        15: {
            "参与人员：": "参与人员：\r音效师：郑青\r音乐设计：郑青 / 朱静雅 / 虞鹏\r音频策划：郑青\r音频技术：冯博宇，蔡新星"
        },
        16: {"资源制作类：": "资源制作类：\r" + "\r".join(u_lines[:6])},
        18: {
            "参与人员": "参与人员\r音效设计：冯博宇\r音乐设计：冯博宇\r音频策划：冯博宇\r音频技术：冯博宇，蔡新星"
        },
        19: {"设计&沟通类：": "设计&沟通类：\r1. OB24音效制作\r2. 多语言分支维护\r3. FM项目音频看板更新"},
        20: {"资源制作类：": "资源制作类：\r" + numbered(fm_lines or ["技能：18个"], 3)},
        22: {
            "参与人员": "参与人员\r音效设计：虞鹏\r音乐设计：虞鹏\r音频策划：虞鹏\r音频技术：蔡新星"
        },
        23: {"设计&沟通类：": "设计&沟通类：\r1. UI替换音频跟进\r2. Stamina功能音频验证\r3. 成就系统音频联调"},
        24: {"资源制作类：": "资源制作类：\r1. 0610 UI替换.mp4\r2. Stamina.mp4\r3. 成就.mp4"},
        27: {
            "平台：": "平台：" + "；".join(tech_cai[3:5]) + "。\rLL：" + "；".join(tech_cai[1:3]) + "。\rU：OB69版本音频功能接入、TW版本功能单。\rL：0618、0716版本对接，语音标准化。\r其他：" + "；".join(tech_zheng[:3]) + "。"
        },
        30: {
            "团队建设": "工具与流程沉淀\r1. 联调工具：实时显示特效、Spine动画和Timeline，点击可定位资源。\r2. Spine查错工具：遍历探索Spine音频事件及父级衰减配置。\r3. External Source编辑器：按标签提取资源并编辑音量。\r4. " + (tech_feng[-1] if tech_feng else "FM项目音频看板更新")
        },
    }


def pptx_shape_type_value(shape) -> int | None:
    try:
        return int(getattr(shape, "shape_type", 0))
    except Exception:
        return None


def is_pptx_media_shape(shape) -> bool:
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.MEDIA or pptx_shape_type_value(shape) == 16


def pptx_shape_rel_ids(shape) -> set[str]:
    rel_ids = set()
    for element in shape.element.iter():
        for value in element.attrib.values():
            if isinstance(value, str) and value.startswith("rId"):
                rel_ids.add(value)
    return rel_ids


def remove_pptx_shape_with_rels(slide, shape):
    for rel_id in pptx_shape_rel_ids(shape):
        try:
            slide.part.drop_rel(rel_id)
        except Exception:
            try:
                del slide.part.rels._rels[rel_id]
            except Exception:
                pass
    parent = shape.element.getparent()
    if parent is not None:
        parent.remove(shape.element)


def prepare_pptx_media_slots(prs: Presentation) -> dict[int, dict[str, float]]:
    slots: dict[int, dict[str, float]] = {}
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in list(slide.shapes):
            if is_pptx_media_shape(shape):
                slots.setdefault(
                    slide_index,
                    {
                        "left": round(int(shape.left) / EMU_PER_POINT, 2),
                        "top": round(int(shape.top) / EMU_PER_POINT, 2),
                        "width": round(int(shape.width) / EMU_PER_POINT, 2),
                        "height": round(int(shape.height) / EMU_PER_POINT, 2),
                    },
                )
                remove_pptx_shape_with_rels(slide, shape)
    return slots


def first_run_font(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                color = None
                try:
                    color = run.font.color.rgb
                except Exception:
                    color = None
                return {
                    "name": run.font.name,
                    "size": run.font.size,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": color,
                }
    return {"name": None, "size": None, "bold": None, "italic": None, "color": None}


def set_pptx_shape_text(shape, text: str):
    style = first_run_font(shape)
    alignment = None
    try:
        alignment = shape.text_frame.paragraphs[0].alignment
    except Exception:
        pass
    lines = str(text).replace("\r", "\n").split("\n")
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
        paragraph.text = line
        if alignment is not None:
            paragraph.alignment = alignment
        for run in paragraph.runs:
            if style.get("name"):
                run.font.name = style["name"]
            if style.get("size"):
                run.font.size = style["size"]
            if style.get("bold") is not None:
                run.font.bold = style["bold"]
            if style.get("italic") is not None:
                run.font.italic = style["italic"]
            if style.get("color") is not None:
                try:
                    run.font.color.rgb = style["color"]
                except Exception:
                    pass


def replace_texts_pptx(prs: Presentation, replacements: dict[int, dict[str, str]]) -> list[dict[str, str]]:
    changed = []
    for slide_index, rules in replacements.items():
        if slide_index > len(prs.slides):
            continue
        slide = prs.slides[slide_index - 1]
        for marker, new_text in rules.items():
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                old_text = getattr(shape, "text", "")
                if marker in old_text:
                    set_pptx_shape_text(shape, new_text)
                    changed.append({"slide": slide_index, "marker": marker, "text": new_text})
                    break
    return changed


def is_av_package_name(name: str) -> bool:
    return Path(str(name).split("#", 1)[0].split("?", 1)[0]).suffix.lower() in MEDIA_EXTS


def remove_transition_sound_xml(pptx_path: Path):
    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    tmp_path = Path(tmp_name)
    try:
        with ZipFile(pptx_path, "r") as source_zip, ZipFile(tmp_path, "w", ZIP_DEFLATED) as output_zip:
            for item in source_zip.infolist():
                if item.filename.startswith("ppt/media/") and is_av_package_name(item.filename):
                    continue
                data = source_zip.read(item.filename)
                if item.filename.endswith(".rels"):
                    root = ET.fromstring(data)
                    changed = False
                    for rel in list(root):
                        target = rel.attrib.get("Target", "")
                        rel_type = rel.attrib.get("Type", "")
                        if ("media/" in target and is_av_package_name(target)) or rel_type.endswith(("/audio", "/video", "/media")):
                            root.remove(rel)
                            changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                elif item.filename.endswith(".xml") and (b"sndAc" in data or b"cMediaNode" in data or b"mediacall" in data):
                    root = ET.fromstring(data)
                    parent_map = {child: parent for parent in root.iter() for child in list(parent)}
                    changed = False
                    for element in list(root.iter()):
                        if element.tag.rsplit("}", 1)[-1] in {"sndAc", "timing"}:
                            parent = parent_map.get(element)
                            if parent is not None:
                                parent.remove(element)
                                changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                elif item.filename == "[Content_Types].xml":
                    root = ET.fromstring(data)
                    changed = False
                    for child in list(root):
                        part_name = child.attrib.get("PartName", "")
                        extension = child.attrib.get("Extension", "")
                        if (part_name.startswith("/ppt/media/") and is_av_package_name(part_name)) or (
                            extension and f".{extension.lower()}" in MEDIA_EXTS
                        ):
                            root.remove(child)
                            changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                output_zip.writestr(item, data)
        shutil.move(str(tmp_path), str(pptx_path))
    finally:
        if tmp_path.exists():
            try:
                tmp_path.unlink()
            except Exception:
                pass


def prepare_staged_pptx(staged_pptx: Path) -> tuple[dict[int, dict[str, float]], list[dict[str, str]]]:
    prs = Presentation(str(LEARN_PPT))
    slots = prepare_pptx_media_slots(prs)
    text_changes = replace_texts_pptx(prs, source_text_plan())
    prs.save(str(staged_pptx))
    remove_transition_sound_xml(staged_pptx)
    return slots, text_changes


def text_shapes(slide) -> list:
    shapes = []
    for index in range(1, slide.Shapes.Count + 1):
        shape = slide.Shapes.Item(index)
        try:
            if shape.HasTextFrame and shape.TextFrame.HasText:
                text = shape.TextFrame.TextRange.Text
                if text and text.strip():
                    shapes.append(shape)
        except Exception:
            continue
    return sorted(shapes, key=lambda item: (float(item.Top), float(item.Left)))


def snapshot_font(shape):
    try:
        font = shape.TextFrame.TextRange.Font
        return {
            "Name": font.Name,
            "NameFarEast": getattr(font, "NameFarEast", font.Name),
            "Size": font.Size,
            "Bold": font.Bold,
            "Color": font.Color.RGB,
        }
    except Exception:
        return {}


def restore_font(shape, font_state):
    try:
        font = shape.TextFrame.TextRange.Font
        if font_state.get("Name"):
            font.Name = font_state["Name"]
        if font_state.get("NameFarEast"):
            font.NameFarEast = font_state["NameFarEast"]
        if font_state.get("Size"):
            font.Size = font_state["Size"]
        if "Bold" in font_state:
            font.Bold = font_state["Bold"]
        if font_state.get("Color") is not None:
            font.Color.RGB = font_state["Color"]
    except Exception:
        pass


def set_shape_text(shape, text: str):
    font_state = snapshot_font(shape)
    shape.TextFrame.TextRange.Text = text
    restore_font(shape, font_state)


def replace_texts(presentation, replacements: dict[int, dict[str, str]]) -> list[dict[str, str]]:
    changed = []
    for slide_index, rules in replacements.items():
        slide = presentation.Slides(slide_index)
        for marker, new_text in rules.items():
            for shape in text_shapes(slide):
                old_text = shape.TextFrame.TextRange.Text
                if marker in old_text:
                    set_shape_text(shape, new_text)
                    changed.append({"slide": slide_index, "marker": marker, "text": new_text})
                    break
    return changed


def is_media_shape(shape) -> bool:
    try:
        if int(shape.Type) == 16:
            return True
    except Exception:
        pass
    try:
        _ = shape.MediaFormat
        return True
    except Exception:
        return False


def clear_media_and_capture_slots(presentation) -> dict[int, dict[str, float]]:
    slots = {}
    for slide_index in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides.Item(slide_index)
        for index in range(slide.Shapes.Count, 0, -1):
            shape = slide.Shapes.Item(index)
            if is_media_shape(shape):
                slots.setdefault(
                    slide_index,
                    {
                        "left": float(shape.Left),
                        "top": float(shape.Top),
                        "width": float(shape.Width),
                        "height": float(shape.Height),
                    },
                )
                shape.Delete()
    return slots


def delete_transition_sounds(presentation):
    for slide_index in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides.Item(slide_index)
        try:
            slide.SlideShowTransition.SoundEffect.Name = ""
        except Exception:
            pass


def ensure_media_slides(presentation, base_indices: list[int], paths: list[Path]) -> list:
    base_slides = [presentation.Slides(index) for index in base_indices]
    if not paths:
        for slide in reversed(base_slides):
            slide.Delete()
        return []
    slots = list(base_slides)
    source = slots[-1]
    after_slide = slots[-1]
    while len(slots) < len(paths):
        duplicate_range = source.Duplicate()
        duplicate_slide = duplicate_range.Item(1)
        duplicate_slide.MoveTo(after_slide.SlideIndex + 1)
        slots.append(duplicate_slide)
        after_slide = duplicate_slide
    for slide in reversed(slots[len(paths) :]):
        slide.Delete()
    return slots[: len(paths)]


def find_title_shape(slide):
    candidates = [shape for shape in text_shapes(slide) if "产出展示" in shape.TextFrame.TextRange.Text]
    if candidates:
        return candidates[0]
    shapes = text_shapes(slide)
    return shapes[0] if shapes else None


def display_title(project_title: str, path: Path) -> str:
    label = "音频展示" if media_kind(path) == "audio" else "产出展示"
    return f"{project_title} {label}：{path.stem}"


def insert_media(slide, path: Path, slot: dict[str, float], project_title: str, embed_limit_mb: int = 80) -> dict[str, str | int | float | bool]:
    title = find_title_shape(slide)
    if title is not None:
        set_shape_text(title, display_title(project_title, path))

    if media_kind(path) == "audio":
        left = 430.0
        top = 240.0
        width = 100.0
        height = 100.0
    else:
        left = float(slot.get("left", 0))
        top = float(slot.get("top", 0))
        width = float(slot.get("width", 960))
        height = float(slot.get("height", 540))

    size_mb = path.stat().st_size / 1024 / 1024
    embed = size_mb <= embed_limit_mb
    # Linked media is still a real PowerPoint media object, but avoids multi-hundred-MB freezes.
    shape = slide.Shapes.AddMediaObject2(str(path), not embed, embed, left, top, width, height)
    shape.Name = path.stem[:64]
    try:
        shape.ZOrder(1)
    except Exception:
        pass
    for text_shape in text_shapes(slide):
        try:
            text_shape.ZOrder(0)
        except Exception:
            pass
    return {
        "slide": int(slide.SlideIndex),
        "path": str(path),
        "name": path.name,
        "kind": media_kind(path),
        "embedded": embed,
        "size_mb": round(size_mb, 2),
        "left": round(left, 2),
        "top": round(top, 2),
        "width": round(width, 2),
        "height": round(height, 2),
    }


def open_presentation(app, path: Path, timeout_seconds: int = 120):
    opened = None
    try:
        opened = app.Presentations.Open(str(path), WithWindow=False)
    except Exception:
        opened = None

    deadline = time.time() + timeout_seconds
    last_error = None
    while time.time() < deadline:
        candidates = []
        if opened is not None:
            candidates.append(opened)
        try:
            count = app.Presentations.Count
            if count:
                candidates.append(app.Presentations.Item(count))
        except Exception as exc:
            last_error = exc

        for candidate in candidates:
            try:
                _ = candidate.Slides.Count
                return candidate
            except Exception as exc:
                last_error = exc
        time.sleep(1)
    raise RuntimeError(f"PowerPoint 打开 PPT 超时：{path} / {last_error}")


def build_with_powerpoint() -> tuple[Path, dict]:
    import win32com.client

    groups = media_groups()
    group_specs = [
        ("TA", [28, 29], "TA"),
        ("EF", [25, 26], "重点工作5——EF项目"),
        ("FM", [21], "重点工作4——FM项目"),
        ("U", [17], "重点工作3——U项目"),
        ("LL", [11, 12, 13, 14], "重点工作2——LL项目"),
        ("L", [6, 7], "重点工作1——L项目"),
    ]

    temp_ppt = Path(tempfile.gettempdir()) / f"codex_pair_stage_{STAMP}.pptx"
    if temp_ppt.exists():
        temp_ppt.unlink()
    slots, text_changes = prepare_staged_pptx(temp_ppt)

    app = None
    presentation = None
    plan: dict = {
        "learn_dir": str(LEARN_DIR),
        "source_dir": str(SOURCE_DIR),
        "learn_ppt": str(LEARN_PPT),
        "strategy": "pair_learned_powerpoint_com_real_media",
        "media_groups": {key: [str(path) for path in value] for key, value in groups.items()},
    }
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            app.DisplayAlerts = 0
        except Exception:
            pass
        app.Visible = 1
        presentation = open_presentation(app, temp_ppt)

        delete_transition_sounds(presentation)

        insertions = []
        for key, base_indices, title in group_specs:
            slides = ensure_media_slides(presentation, base_indices, groups[key])
            fallback_slot = slots.get(base_indices[-1], {"left": 0, "top": 0, "width": 960, "height": 540})
            for offset, (slide, path) in enumerate(zip(slides, groups[key])):
                original_slot = slots.get(base_indices[min(offset, len(base_indices) - 1)], fallback_slot)
                insertions.append(insert_media(slide, path, original_slot, title))

        plan["text_changes"] = text_changes
        plan["media_insertions"] = insertions
        plan["captured_media_slots"] = slots
        presentation.SaveAs(str(PPTX_OUT))
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        time.sleep(1)
        try:
            temp_ppt.unlink()
        except Exception:
            pass

    PLAN_OUT.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")
    return PPTX_OUT, plan


def main() -> int:
    output, plan = build_with_powerpoint()
    print(f"PPTX={output}")
    print(f"PLAN={PLAN_OUT}")
    print(f"MEDIA_INSERTIONS={len(plan.get('media_insertions', []))}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
